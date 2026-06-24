from __future__ import annotations

import re
from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

from .constants import A_NS, BLACK, BLANK_LINE, MATH_FONT, WHITE
from .design import get_design_value, rgb_from_hex

SUPERSCRIPT_TO_NORMAL = {
    "⁰": "0",
    "¹": "1",
    "²": "2",
    "³": "3",
    "⁴": "4",
    "⁵": "5",
    "⁶": "6",
    "⁷": "7",
    "⁸": "8",
    "⁹": "9",
}

NORMAL_TO_SUPERSCRIPT = {value: key for key, value in SUPERSCRIPT_TO_NORMAL.items()}
SUPERSCRIPT_DIGITS = "".join(SUPERSCRIPT_TO_NORMAL)
SUPERSCRIPT_DIGIT_RE = f"[{SUPERSCRIPT_DIGITS}]+"

FRACTION_RE = re.compile(
    rf"(?P<num>-?\([^)]*\)|-?[0-9A-Za-z가-힣{SUPERSCRIPT_DIGITS}]+)\s*/\s*"
    rf"(?P<den>-?\([^)]*\)|-?\[\s*\](?:{SUPERSCRIPT_DIGIT_RE}|\^\d+)?|-?[0-9A-Za-z가-힣{SUPERSCRIPT_DIGITS}×]+)"
)
BLANK_RE = re.compile(r"\^\[\s*\]|\[\s*\]|__|□")
DENOMINATOR_BLANK_EXP_RE = re.compile(rf"\[\s*\](?P<exp>{SUPERSCRIPT_DIGIT_RE})?")
INLINE_OBJECT_RE = re.compile(r"\(\s*\)|[=÷×]")
INLINE_OPERATOR_TOKENS = {"=", "÷", "×"}
INLINE_NUMERIC_TOKEN_RE = re.compile(r"-?[0-9][0-9.,…]*")
ANSWER_PAREN_TOKEN = "( )"
XML_SPACE = "{http://www.w3.org/XML/1998/namespace}space"

THEME_COLORS = {
    "BACKGROUND_1": MSO_THEME_COLOR.BACKGROUND_1,
    "BACKGROUND_2": MSO_THEME_COLOR.BACKGROUND_2,
    "TEXT_1": MSO_THEME_COLOR.TEXT_1,
    "TEXT_2": MSO_THEME_COLOR.TEXT_2,
    "ACCENT_1": MSO_THEME_COLOR.ACCENT_1,
    "ACCENT_2": MSO_THEME_COLOR.ACCENT_2,
    "ACCENT_3": MSO_THEME_COLOR.ACCENT_3,
    "ACCENT_4": MSO_THEME_COLOR.ACCENT_4,
    "ACCENT_5": MSO_THEME_COLOR.ACCENT_5,
    "ACCENT_6": MSO_THEME_COLOR.ACCENT_6,
}


def is_bt_math_font(font: str) -> bool:
    return font.strip().casefold() == MATH_FONT.casefold()


def style_value(design: dict[str, Any] | None, path: str, default: Any) -> Any:
    return get_design_value(design, path, default)


def style_color(design: dict[str, Any] | None, name: str, default: RGBColor) -> RGBColor:
    return rgb_from_hex(get_design_value(design, f"colors.{name}", None), default)


def style_theme_color(design: dict[str, Any] | None, path: str) -> MSO_THEME_COLOR | None:
    value = style_value(design, path, None)
    if not value:
        return None
    return THEME_COLORS.get(str(value).strip().upper())


def apply_blank_style(blank, design: dict[str, Any] | None, *, line_scale: float = 1.0) -> None:
    fill_mode = str(style_value(design, "problem_slide.item.blank_fill", "")).strip().lower()
    if fill_mode in {"none", "nofill", "transparent"}:
        blank.fill.background()
    else:
        blank.fill.solid()
        fill_theme = style_theme_color(design, "problem_slide.item.blank_fill_theme_color")
        if fill_theme is not None:
            blank.fill.fore_color.theme_color = fill_theme
        else:
            blank.fill.fore_color.rgb = style_color(design, "blank_fill", WHITE)

    line_theme = style_theme_color(design, "problem_slide.item.blank_line_theme_color")
    if line_theme is not None:
        blank.line.color.theme_color = line_theme
    else:
        blank.line.color.rgb = style_color(design, "blank_line", BLANK_LINE)
    blank.line.width = Pt(style_value(design, "problem_slide.item.blank_line_pt", 2.0) * line_scale)
    apply_blank_line_xml_style(blank, design)


def apply_blank_line_xml_style(blank, design: dict[str, Any] | None) -> None:
    sp_pr = blank.element.find(qn("p:spPr"))
    if sp_pr is None:
        return
    line = sp_pr.find(qn("a:ln"))
    if line is None:
        return

    cap = style_value(design, "problem_slide.item.blank_line_cap", None)
    if cap:
        line.set("cap", str(cap))

    solid_fill = line.find(qn("a:solidFill"))
    scheme_clr = solid_fill.find(qn("a:schemeClr")) if solid_fill is not None else None
    lum_mod = style_value(design, "problem_slide.item.blank_line_lum_mod", None)
    if scheme_clr is not None and lum_mod is not None:
        for child in list(scheme_clr):
            if child.tag == qn("a:lumMod"):
                scheme_clr.remove(child)
        lum = OxmlElement("a:lumMod")
        lum.set("val", str(int(lum_mod)))
        scheme_clr.append(lum)

    join = str(style_value(design, "problem_slide.item.blank_line_join", "")).strip().lower()
    if join == "round" and line.find(qn("a:round")) is None:
        line.append(OxmlElement("a:round"))


def normalize_caret_exponents(text: str) -> str:
    def replace(match: re.Match[str]) -> str:
        return "".join(NORMAL_TO_SUPERSCRIPT.get(ch, ch) for ch in match.group(1))

    return re.sub(r"\^(\d+)", replace, text)


def rich_text_segments(text: str) -> list[tuple[str, str]]:
    text = normalize_caret_exponents(text)
    segments: list[tuple[str, str]] = []
    buffer: list[str] = []

    for char in text:
        if char in SUPERSCRIPT_TO_NORMAL:
            if buffer:
                segments.append(("normal", "".join(buffer)))
                buffer = []
            segments.append(("superscript", SUPERSCRIPT_TO_NORMAL[char]))
        else:
            buffer.append(char)

    if buffer:
        segments.append(("normal", "".join(buffer)))
    return segments


def is_answer_parenthesis_token(token: str) -> bool:
    return re.fullmatch(r"\(\s*\)", token) is not None


def is_inline_numeric_token(token: str) -> bool:
    return INLINE_NUMERIC_TOKEN_RE.fullmatch(normalize_caret_exponents(token).strip()) is not None


def preserve_repeated_spaces(run, value: str) -> None:
    if "  " not in value:
        return
    text_node = run._r.find(f"{{{A_NS}}}t")
    if text_node is not None:
        text_node.set(XML_SPACE, "preserve")


def estimate_text_width(text: str, size: float) -> float:
    scale = size / 21.0
    width = 0.0
    for char in text:
        if char.isspace():
            width += 0.055
        elif char in ".,:;()[]":
            width += 0.055
        elif char in "=+-×÷":
            width += 0.12
        elif char in SUPERSCRIPT_TO_NORMAL:
            width += 0.065
        elif ord(char) > 127:
            width += 0.13
        else:
            width += 0.105
    return max(0.06, width * scale)


def clear_text_frame(shape) -> None:
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = False
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    text_frame.margin_left = Inches(0)
    text_frame.margin_right = Inches(0)
    text_frame.margin_top = Inches(0)
    text_frame.margin_bottom = Inches(0)
    text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT


def add_rich_text(
    container,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    *,
    size: float,
    font: str = MATH_FONT,
    color: RGBColor = BLACK,
    align=PP_ALIGN.LEFT,
) -> None:
    shape = container.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    clear_text_frame(shape)
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.alignment = align

    for kind, value in rich_text_segments(text):
        run = paragraph.add_run()
        run.text = value
        preserve_repeated_spaces(run, value)
        if kind == "superscript":
            run.font.size = Pt(size * 0.75)
            run._r.get_or_add_rPr().set("baseline", "60000")
        else:
            run.font.size = Pt(size)
        run.font.name = font
        run.font.bold = is_bt_math_font(font)
        run.font.color.rgb = color


def add_blank_shape(
    container,
    x: float,
    y: float,
    design: dict[str, Any] | None,
    *,
    scale: float = 1.0,
    line_scale: float = 1.0,
) -> float:
    width = style_value(design, "problem_slide.item.blank_w", 0.55) * scale
    height = style_value(design, "problem_slide.item.blank_h", 0.42) * scale
    blank = container.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height))
    apply_blank_style(blank, design, line_scale=line_scale)
    return width


def is_explicit_exponent_blank(token: str) -> bool:
    return token.startswith("^")


def is_exponent_blank_match(text: str, match: re.Match[str], design: dict[str, Any] | None) -> bool:
    return is_explicit_exponent_blank(match.group(0))


def exponent_blank_width(design: dict[str, Any] | None) -> float:
    return style_value(design, "problem_slide.item.blank_w", 0.55) * style_value(
        design,
        "math.exponent_blank_scale",
        0.48,
    )


def add_exponent_blank_shape(
    container,
    x: float,
    y: float,
    design: dict[str, Any] | None,
) -> float:
    return add_blank_shape(
        container,
        x,
        y + style_value(design, "math.exponent_blank_y_offset", -0.03),
        design,
        scale=style_value(design, "math.exponent_blank_scale", 0.48),
        line_scale=style_value(design, "math.exponent_blank_line_scale", 0.55),
    )


def measure_inline_segment_width(
    text: str,
    *,
    design: dict[str, Any] | None,
    size: float,
    max_width: float,
    blank_scale: float = 1.0,
) -> float:
    width = 0.0
    text = normalize_caret_exponents(text)
    pos = 0
    for match in BLANK_RE.finditer(text):
        prefix = text[pos : match.start()]
        if prefix:
            width += min(measure_inline_text_width(prefix, design=design, size=size), max_width - width)
        if width < max_width:
            if is_exponent_blank_match(text, match, design):
                width += style_value(design, "math.exponent_blank_pre_gap", 0.16)
                width += exponent_blank_width(design)
                width += style_value(design, "math.exponent_blank_gap", 0.04)
            else:
                width += style_value(design, "math.inline_blank_pre_gap", 0.0)
                width += style_value(design, "problem_slide.item.blank_w", 0.55) * blank_scale
                width += style_value(design, "math.inline_blank_post_gap", 0.08)
        pos = match.end()

    suffix = text[pos:]
    if suffix and width < max_width:
        width += min(measure_inline_text_width(suffix, design=design, size=size), max_width - width)
    return min(width, max_width)


def inline_text_tokens(text: str, design: dict[str, Any] | None) -> list[str]:
    if not style_value(design, "math.split_inline_operator_objects", True):
        stripped = text.strip()
        return [stripped] if stripped else []

    tokens: list[str] = []
    pos = 0
    for match in INLINE_OBJECT_RE.finditer(text):
        prefix = text[pos : match.start()].strip()
        if prefix:
            tokens.append(prefix)
        matched = match.group(0)
        if matched.startswith("("):
            tokens.append(
                matched
                if style_value(design, "math.preserve_answer_parenthesis_spaces", True)
                else ANSWER_PAREN_TOKEN
            )
        else:
            tokens.append(matched)
        pos = match.end()

    suffix = text[pos:].strip()
    if suffix:
        tokens.append(suffix)
    return tokens


def answer_parenthesis_token_width(token: str, size: float, design: dict[str, Any] | None) -> float:
    char_em = style_value(design, "math.answer_parenthesis_char_em", None)
    padding = style_value(design, "math.answer_parenthesis_width_padding", 0.0)
    min_width = style_value(design, "math.answer_parenthesis_min_w", 0.0)
    if char_em is not None:
        char_slot_width = (size / 72.0) * char_em
        return max(char_slot_width * len(token) + padding, min_width)

    base_width = estimate_text_width(token, size)
    scale = style_value(design, "math.answer_parenthesis_width_scale", 2.0)
    return max(base_width * scale + padding, min_width)


def inline_token_width(token: str, size: float, design: dict[str, Any] | None = None) -> float:
    if token in INLINE_OPERATOR_TOKENS:
        return estimate_text_width(f" {token} ", size)
    if is_answer_parenthesis_token(token):
        return answer_parenthesis_token_width(token, size, design)
    width = estimate_text_width(token, size)
    if is_inline_numeric_token(token):
        scale = style_value(design, "math.inline_numeric_width_scale", 1.0)
        padding = style_value(design, "math.inline_numeric_width_padding", 0.0)
        min_width = style_value(design, "math.inline_numeric_min_w", 0.0)
        width = max(width * scale + padding, min_width)
    else:
        width *= style_value(design, "math.inline_math_width_scale", 1.35)
    return width


def measure_inline_text_width(text: str, *, design: dict[str, Any] | None, size: float) -> float:
    return sum(inline_token_width(token, size, design) for token in inline_text_tokens(text, design))


def inline_segment_text_alignment(design: dict[str, Any] | None, requested_align, token: str) -> Any:
    if requested_align == PP_ALIGN.CENTER:
        return PP_ALIGN.CENTER
    if (
        style_value(design, "math.center_inline_operator_text", True)
        and (token in INLINE_OPERATOR_TOKENS or is_answer_parenthesis_token(token))
    ):
        return PP_ALIGN.CENTER
    return PP_ALIGN.LEFT


def add_inline_text(
    container,
    x: float,
    y: float,
    text: str,
    *,
    design: dict[str, Any] | None,
    size: float,
    font: str,
    max_width: float,
    align=PP_ALIGN.LEFT,
) -> float:
    cursor = x
    for token in inline_text_tokens(text, design):
        if cursor - x >= max_width:
            break
        width = min(inline_token_width(token, size, design), max_width - (cursor - x))
        add_rich_text(
            container,
            cursor,
            y,
            width,
            0.42,
            token,
            size=size,
            font=font,
            align=inline_segment_text_alignment(design, align, token),
        )
        cursor += width
    return cursor - x


def add_inline_segment(
    container,
    x: float,
    y: float,
    text: str,
    *,
    design: dict[str, Any] | None,
    size: float,
    font: str,
    max_width: float,
    blank_scale: float = 1.0,
    align=PP_ALIGN.LEFT,
) -> float:
    cursor = x
    text = normalize_caret_exponents(text)
    if align == PP_ALIGN.CENTER:
        used_width = measure_inline_segment_width(
            text,
            design=design,
            size=size,
            max_width=max_width,
            blank_scale=blank_scale,
        )
        cursor += max(0.0, (max_width - used_width) / 2)
    pos = 0
    for match in BLANK_RE.finditer(text):
        prefix = text[pos : match.start()]
        if prefix:
            cursor += add_inline_text(
                container,
                cursor,
                y,
                prefix,
                design=design,
                size=size,
                font=font,
                max_width=max_width - (cursor - x),
                align=align,
            )
        if cursor - x < max_width:
            if is_exponent_blank_match(text, match, design):
                cursor += style_value(design, "math.exponent_blank_pre_gap", 0.16)
                cursor += add_exponent_blank_shape(container, cursor, y, design)
                cursor += style_value(design, "math.exponent_blank_gap", 0.04)
            else:
                cursor += style_value(design, "math.inline_blank_pre_gap", 0.0)
                cursor += add_blank_shape(container, cursor, y + 0.03, design, scale=blank_scale)
                cursor += style_value(design, "math.inline_blank_post_gap", 0.08)
        pos = match.end()

    suffix = text[pos:]
    if suffix and cursor - x < max_width:
        cursor += add_inline_text(
            container,
            cursor,
            y,
            suffix,
            design=design,
            size=size,
            font=font,
            max_width=max_width - (cursor - x),
            align=align,
        )
    return cursor - x


def cleaned_fraction_text(text: str) -> str:
    return text.strip()


def is_multiplicative_group(text: str) -> bool:
    return re.fullmatch(r"[\[\]\s0-9A-Za-z가-힣⁰¹²³⁴⁵⁶⁷⁸⁹×^]+", text) is not None


def cleaned_denominator_text(text: str) -> str:
    text = cleaned_fraction_text(text)
    if text.startswith("(") and text.endswith(")"):
        inner = text[1:-1].strip()
        if inner and is_multiplicative_group(inner):
            return inner
    return text


def normal_digits(text: str) -> str:
    return "".join(SUPERSCRIPT_TO_NORMAL.get(char, char) for char in text)


def denominator_blank_exponent_tokens(text: str) -> list[tuple[str, str | None]] | None:
    text = normalize_caret_exponents(cleaned_fraction_text(text))
    tokens: list[tuple[str, str | None]] = []
    has_blank_exponent = False
    pos = 0
    while pos < len(text):
        if text[pos].isspace():
            pos += 1
            continue
        blank_match = DENOMINATOR_BLANK_EXP_RE.match(text, pos)
        if blank_match:
            exponent = blank_match.group("exp")
            tokens.append(("blank", exponent))
            has_blank_exponent = has_blank_exponent or bool(exponent)
            pos = blank_match.end()
            continue
        if text[pos] == "×":
            tokens.append(("operator", "×"))
            pos += 1
            continue
        return None
    return tokens if has_blank_exponent else None


def fraction_blank_exponent_text_width(text: str, design: dict[str, Any] | None) -> float | None:
    tokens = denominator_blank_exponent_tokens(text)
    if tokens is None:
        return None
    blank_w = style_value(design, "problem_slide.item.blank_w", 0.55) * style_value(
        design,
        "math.fraction_blank_scale",
        1.0,
    )
    exp_gap = style_value(design, "math.fraction_den_blank_exp_gap", 0.023)
    exp_w = style_value(design, "math.fraction_den_blank_exp_w", 0.133)
    operator_w = style_value(design, "math.fraction_den_operator_w", 0.345)
    blank_operator_gap = style_value(design, "math.fraction_den_blank_operator_gap", 0.111)
    operator_blank_gap = style_value(design, "math.fraction_den_operator_blank_gap", 0.048)
    cursor = 0.0
    extent = 0.0
    for kind, value in tokens:
        if kind == "blank":
            extent = max(extent, cursor + blank_w)
            if value:
                extent = max(extent, cursor + blank_w + exp_gap + exp_w)
            cursor += blank_w
            continue
        if kind == "operator":
            cursor += blank_operator_gap
            extent = max(extent, cursor + operator_w)
            cursor += operator_w + operator_blank_gap
    return (
        style_value(design, "math.fraction_den_blank_exp_left_padding", 0.08)
        + max(extent, cursor)
        + style_value(design, "math.fraction_den_blank_exp_right_padding", 0.04)
    )


def fraction_plain_text_width(text: str, design: dict[str, Any] | None, size: float) -> float:
    text = normalize_caret_exponents(cleaned_fraction_text(text))
    visible = [char for char in text if not char.isspace()]
    if not visible:
        return style_value(design, "math.fraction_min_w", 0.355)

    if any(char in "×÷/" or char in SUPERSCRIPT_TO_NORMAL for char in visible):
        units = 0.0
        for char in visible:
            if char in SUPERSCRIPT_TO_NORMAL:
                units += style_value(design, "math.fraction_exponent_unit", 0.75)
            elif char in "×÷/":
                units += style_value(design, "math.fraction_operator_unit", 1.4)
            elif char in "()":
                units += 0.25
            else:
                units += 1.0
        return (
            units * style_value(design, "math.fraction_multiplicative_char_w", 0.2) * (size / 24.0)
            + style_value(design, "math.fraction_multiplicative_padding", 0.1)
        )

    return (
        len(visible) * style_value(design, "math.fraction_numeric_char_w", 0.157) * (size / 24.0)
        + style_value(design, "math.fraction_numeric_padding", 0.158)
    )


def fraction_blank_text_width(text: str, design: dict[str, Any] | None, size: float) -> float:
    blank_count = len(BLANK_RE.findall(text))
    blank_scale = style_value(design, "math.fraction_blank_scale", 1.0)
    base = measure_inline_segment_width(
        text,
        design=design,
        size=size * 0.86,
        max_width=99.0,
        blank_scale=blank_scale,
    )
    return (
        base
        + style_value(design, "math.fraction_blank_width_padding", 0.18)
        + max(0, blank_count - 1) * style_value(design, "math.fraction_extra_blank_width_padding", 0.47)
    )


def fraction_text_width(text: str, design: dict[str, Any] | None, size: float) -> float:
    text = cleaned_fraction_text(text)
    blank_exponent_width = fraction_blank_exponent_text_width(text, design)
    if blank_exponent_width is not None:
        return blank_exponent_width
    if BLANK_RE.search(text):
        return fraction_blank_text_width(text, design, size)
    return fraction_plain_text_width(text, design, size)


def fraction_denominator_y(
    y: float,
    *,
    design: dict[str, Any] | None,
    num_y: float,
    bar_y: float,
    num_h: float,
    bar_h: float,
) -> float:
    if not style_value(design, "math.fraction_equal_vertical_gap", False):
        return y + style_value(design, "math.fraction_den_y_offset", 0.25)
    numerator_to_bar_gap = max(0.0, bar_y - (num_y + num_h))
    return bar_y + bar_h + numerator_to_bar_gap


def add_fraction(
    container,
    x: float,
    y: float,
    numerator: str,
    denominator: str,
    *,
    design: dict[str, Any] | None,
    size: float,
    font: str,
) -> float:
    numerator = cleaned_fraction_text(numerator)
    denominator = cleaned_denominator_text(denominator)
    fraction_width = max(
        style_value(design, "math.fraction_min_w", 0.42),
        fraction_text_width(numerator, design, size),
        fraction_text_width(denominator, design, size),
    )
    num_h = style_value(design, "math.fraction_text_h", 0.28)
    den_h = style_value(design, "math.fraction_text_h", 0.28)
    bar_h = style_value(design, "math.fraction_bar_h", 0.018)
    num_y = y + style_value(design, "math.fraction_num_y_offset", -0.23)
    bar_y = y + style_value(design, "math.fraction_bar_y_offset", 0.24)
    den_y = fraction_denominator_y(
        y,
        design=design,
        num_y=num_y,
        bar_y=bar_y,
        num_h=num_h,
        bar_h=bar_h,
    )

    add_rich_text(
        container,
        x,
        num_y,
        fraction_width,
        num_h,
        numerator,
        size=size,
        font=font,
        align=PP_ALIGN.CENTER,
    )
    bar = container.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(bar_y), Inches(fraction_width), Inches(bar_h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLACK
    bar.line.fill.background()

    if denominator_blank_exponent_tokens(denominator) is not None:
        add_fraction_blank_exponent_denominator(
            container,
            x,
            bar_y,
            denominator,
            fraction_width=fraction_width,
            design=design,
            size=size,
            font=font,
        )
    elif BLANK_RE.search(denominator):
        blank_scale = style_value(design, "math.fraction_blank_scale", 1.0)
        add_inline_segment(
            container,
            x,
            den_y,
            denominator,
            design=design,
            size=size * 0.86,
            font=font,
            max_width=fraction_width,
            blank_scale=blank_scale,
            align=PP_ALIGN.CENTER,
        )
    else:
        add_rich_text(
            container,
            x,
            den_y,
            fraction_width,
            den_h,
            denominator,
            size=size,
            font=font,
            align=PP_ALIGN.CENTER,
        )
    return fraction_width


def add_fraction_blank_exponent_denominator(
    container,
    x: float,
    bar_y: float,
    denominator: str,
    *,
    fraction_width: float,
    design: dict[str, Any] | None,
    size: float,
    font: str,
) -> None:
    tokens = denominator_blank_exponent_tokens(denominator)
    if tokens is None:
        return
    content_width = fraction_blank_exponent_text_width(denominator, design) or fraction_width
    content_x = x + max(0.0, (fraction_width - content_width) / 2)
    cursor = content_x + style_value(design, "math.fraction_den_blank_exp_left_padding", 0.08)
    blank_y = bar_y + style_value(design, "math.fraction_den_blank_bar_gap", 0.089)
    blank_w = style_value(design, "problem_slide.item.blank_w", 0.55) * style_value(
        design,
        "math.fraction_blank_scale",
        1.0,
    )
    exp_gap = style_value(design, "math.fraction_den_blank_exp_gap", 0.023)
    exp_w = style_value(design, "math.fraction_den_blank_exp_w", 0.133)
    exp_h = style_value(design, "math.fraction_den_blank_exp_h", 0.303)
    exp_y = blank_y + style_value(design, "math.fraction_den_blank_exp_y_offset", -0.112)
    exp_size = style_value(design, "math.fraction_den_blank_exp_font_size", size * 0.83)
    operator_w = style_value(design, "math.fraction_den_operator_w", 0.345)
    operator_h = style_value(design, "math.fraction_den_operator_h", 0.404)
    operator_size = style_value(design, "math.fraction_den_operator_font_size", size)
    operator_y = blank_y + style_value(design, "math.fraction_den_operator_y_offset", 0.045)
    blank_operator_gap = style_value(design, "math.fraction_den_blank_operator_gap", 0.111)
    operator_blank_gap = style_value(design, "math.fraction_den_operator_blank_gap", 0.048)

    for kind, value in tokens:
        if kind == "blank":
            add_blank_shape(
                container,
                cursor,
                blank_y,
                design,
                scale=style_value(design, "math.fraction_blank_scale", 1.0),
            )
            if value:
                add_rich_text(
                    container,
                    cursor + blank_w + exp_gap,
                    exp_y,
                    exp_w,
                    exp_h,
                    normal_digits(value),
                    size=exp_size,
                    font=font,
                    align=PP_ALIGN.CENTER,
                )
            cursor += blank_w
            continue
        if kind == "operator":
            cursor += blank_operator_gap
            add_rich_text(
                container,
                cursor,
                operator_y,
                operator_w,
                operator_h,
                value or "",
                size=operator_size,
                font=font,
                align=PP_ALIGN.CENTER,
            )
            cursor += operator_w + operator_blank_gap


def add_math_row(
    container,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    *,
    design: dict[str, Any] | None,
    size: float,
    font: str = MATH_FONT,
) -> float:
    """Render a compact editable math row with simple fractions and exponents.

    This is intentionally conservative: it recognizes slash fractions in the
    common workbook patterns and leaves other text editable as rich text.
    """

    cursor = x
    text = normalize_caret_exponents(text)
    pos = 0
    for match in FRACTION_RE.finditer(text):
        prefix = text[pos : match.start()]
        if prefix:
            cursor += add_inline_segment(
                container,
                cursor,
                y,
                prefix,
                design=design,
                size=size,
                font=font,
                max_width=max(0.0, x + w - cursor),
            )
        if cursor >= x + w:
            return w
        cursor += add_fraction(
            container,
            cursor,
            y,
            match.group("num"),
            match.group("den"),
            design=design,
            size=size,
            font=font,
        )
        cursor += style_value(design, "math.token_gap", 0.08)
        pos = match.end()

    suffix = text[pos:]
    if suffix and cursor < x + w:
        cursor += add_inline_segment(
            container,
            cursor,
            y,
            suffix,
            design=design,
            size=size,
            font=font,
            max_width=max(0.0, x + w - cursor),
        )
    return max(0.0, min(cursor - x, w))
