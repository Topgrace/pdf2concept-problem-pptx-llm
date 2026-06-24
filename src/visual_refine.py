from __future__ import annotations

import json
import re
import shutil
import subprocess
from pathlib import Path
from typing import Any

import fitz
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

try:
    from PIL import Image, ImageDraw, ImageFont
except Exception:  # pragma: no cover - optional runtime fallback
    Image = None
    ImageDraw = None
    ImageFont = None


EMU_PER_INCH = 914400
EXPONENT_BLANK_RE = re.compile(r"\^\[\s*\]")
BLANK_RE = re.compile(r"\^\[\s*\]|\[\s*\]|__|□")
BASELINE_BLANK_RE = re.compile(r"(?<!\^)\[\s*\]|__|□")


def set_nested_value(data: dict[str, Any], dotted_path: str, value: Any) -> None:
    target = data
    parts = dotted_path.split(".")
    for part in parts[:-1]:
        target = target.setdefault(part, {})
    target[parts[-1]] = value


def apply_design_adjustments(design: dict[str, Any], adjustments: dict[str, Any]) -> dict[str, Any]:
    if not adjustments:
        return design
    cloned = json.loads(json.dumps(design, ensure_ascii=False))
    for path, value in adjustments.items():
        set_nested_value(cloned, path, value)
    return cloned


def count_blank_intent(generation_report: dict[str, Any]) -> dict[str, int]:
    exponent = 0
    baseline = 0
    for block in generation_report.get("block_inventory", []):
        for item in block.get("items", []):
            raw_text = str(item.get("raw_text", ""))
            exponent += len(EXPONENT_BLANK_RE.findall(raw_text))
            baseline += len(BASELINE_BLANK_RE.findall(raw_text))
    return {"exponent_blank_count": exponent, "baseline_blank_count": baseline}


def render_pdf_pages(pdf_path: Path, pages: list[int], output_dir: Path, dpi: int = 144) -> list[dict[str, Any]]:
    output_dir.mkdir(parents=True, exist_ok=True)
    doc = fitz.open(pdf_path)
    matrix = fitz.Matrix(dpi / 72, dpi / 72)
    rendered: list[dict[str, Any]] = []
    for page_no in pages:
        page = doc[page_no - 1]
        image_path = output_dir / f"source-page-{page_no:03d}.png"
        pix = page.get_pixmap(matrix=matrix, alpha=False)
        pix.save(image_path)
        rendered.append({"page": page_no, "image": str(image_path), "width_px": pix.width, "height_px": pix.height})
    return rendered


def export_pptx_with_powerpoint(pptx_path: Path, output_dir: Path) -> list[dict[str, Any]]:
    import win32com.client

    output_dir.mkdir(parents=True, exist_ok=True)
    app = win32com.client.DispatchEx("PowerPoint.Application")
    presentation = None
    try:
        presentation = app.Presentations.Open(str(pptx_path.resolve()), WithWindow=False)
        rendered: list[dict[str, Any]] = []
        for index in range(1, presentation.Slides.Count + 1):
            image_path = output_dir / f"slide-{index:03d}.png"
            presentation.Slides(index).Export(str(image_path.resolve()), "PNG")
            rendered.append({"slide_number": index, "image": str(image_path)})
        return rendered
    finally:
        if presentation is not None:
            presentation.Close()
        app.Quit()


def export_pptx_with_soffice(pptx_path: Path, output_dir: Path) -> list[dict[str, Any]]:
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        raise RuntimeError("PowerPoint COM and LibreOffice export are unavailable.")
    output_dir.mkdir(parents=True, exist_ok=True)
    subprocess.run(
        [soffice, "--headless", "--convert-to", "png", "--outdir", str(output_dir), str(pptx_path)],
        check=True,
        capture_output=True,
        text=True,
    )
    return [
        {"slide_number": idx + 1, "image": str(path)}
        for idx, path in enumerate(sorted(output_dir.glob("*.png")))
    ]


def export_pptx_geometry_snapshots(pptx_path: Path, output_dir: Path) -> list[dict[str, Any]]:
    if Image is None or ImageDraw is None:
        raise RuntimeError("No image backend available for PPTX visual snapshots.")
    output_dir.mkdir(parents=True, exist_ok=True)
    prs = Presentation(pptx_path)
    rendered: list[dict[str, Any]] = []
    scale = 120
    width = int(prs.slide_width / EMU_PER_INCH * scale)
    height = int(prs.slide_height / EMU_PER_INCH * scale)
    font = ImageFont.load_default() if ImageFont is not None else None

    def draw_shapes(draw, shapes) -> None:
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                draw_shapes(draw, shape.shapes)
                continue
            x0 = int(shape.left / EMU_PER_INCH * scale)
            y0 = int(shape.top / EMU_PER_INCH * scale)
            x1 = int((shape.left + shape.width) / EMU_PER_INCH * scale)
            y1 = int((shape.top + shape.height) / EMU_PER_INCH * scale)
            try:
                auto_type = str(shape.auto_shape_type).lower()
            except Exception:
                auto_type = ""
            if "round" in auto_type:
                draw.rounded_rectangle([x0, y0, x1, y1], radius=8, outline=(120, 120, 120), width=2, fill=(255, 255, 255))
            elif getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    draw.text((x0, y0), text[:80], fill=(0, 0, 0), font=font)

    for index, slide in enumerate(prs.slides, start=1):
        image = Image.new("RGB", (width, height), "white")
        draw = ImageDraw.Draw(image)
        draw_shapes(draw, slide.shapes)
        image_path = output_dir / f"slide-{index:03d}.png"
        image.save(image_path)
        rendered.append({"slide_number": index, "image": str(image_path), "backend": "geometry_snapshot"})
    return rendered


def export_pptx_pngs(pptx_path: Path, output_dir: Path) -> tuple[list[dict[str, Any]], str, str | None]:
    try:
        return export_pptx_with_powerpoint(pptx_path, output_dir), "powerpoint", None
    except Exception as exc:
        powerpoint_error = str(exc)
    try:
        return export_pptx_with_soffice(pptx_path, output_dir), "libreoffice", powerpoint_error
    except Exception as exc:
        fallback_error = f"PowerPoint: {powerpoint_error}; fallback: {exc}"
    return export_pptx_geometry_snapshots(pptx_path, output_dir), "geometry_snapshot", fallback_error


def make_comparison_images(
    source_images: list[dict[str, Any]],
    slide_images: list[dict[str, Any]],
    output_dir: Path,
) -> list[dict[str, Any]]:
    if Image is None:
        return []
    output_dir.mkdir(parents=True, exist_ok=True)
    comparisons: list[dict[str, Any]] = []
    source_by_index = {idx: item for idx, item in enumerate(source_images)}
    problem_slides = [item for item in slide_images if item.get("slide_number", 0) > 1]
    for idx, slide in enumerate(problem_slides):
        source = source_by_index.get(min(idx, max(source_by_index) if source_by_index else 0))
        if not source:
            continue
        left = Image.open(source["image"]).convert("RGB")
        right = Image.open(slide["image"]).convert("RGB")
        target_h = max(left.height, right.height)
        left = left.resize((int(left.width * target_h / left.height), target_h))
        right = right.resize((int(right.width * target_h / right.height), target_h))
        canvas = Image.new("RGB", (left.width + right.width, target_h), "white")
        canvas.paste(left, (0, 0))
        canvas.paste(right, (left.width, 0))
        output = output_dir / f"compare-slide-{slide['slide_number']:03d}.png"
        canvas.save(output)
        comparisons.append({"slide_number": slide["slide_number"], "image": str(output)})
    return comparisons


def iter_shapes(shapes):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)
        else:
            yield shape


def bbox_in(shape) -> tuple[float, float, float, float]:
    return (
        shape.left / EMU_PER_INCH,
        shape.top / EMU_PER_INCH,
        (shape.left + shape.width) / EMU_PER_INCH,
        (shape.top + shape.height) / EMU_PER_INCH,
    )


def intersection_area(a: tuple[float, float, float, float], b: tuple[float, float, float, float]) -> float:
    x0 = max(a[0], b[0])
    y0 = max(a[1], b[1])
    x1 = min(a[2], b[2])
    y1 = min(a[3], b[3])
    if x1 <= x0 or y1 <= y0:
        return 0.0
    return (x1 - x0) * (y1 - y0)


def inspect_exponent_blank_geometry(
    pptx_path: Path,
    generation_report: dict[str, Any],
    design: dict[str, Any],
) -> dict[str, Any]:
    intent = count_blank_intent(generation_report)
    prs = Presentation(pptx_path)
    problem_slide_numbers = {slide["slide_number"] for slide in generation_report.get("slide_trace", [])}
    normal_w = float(design.get("problem_slide", {}).get("item", {}).get("blank_w", 0.55))
    normal_h = float(design.get("problem_slide", {}).get("item", {}).get("blank_h", 0.551))
    expected_scale = float(design.get("math", {}).get("exponent_blank_scale", 0.48))
    max_exp_w = normal_w * max(0.68, expected_scale + 0.15)
    max_exp_h = normal_h * max(0.68, expected_scale + 0.15)
    min_base_w = normal_w * 0.78
    min_base_h = normal_h * 0.78

    small_blanks: list[dict[str, Any]] = []
    baseline_blanks: list[dict[str, Any]] = []
    overlaps: list[dict[str, Any]] = []
    raised_position_issues: list[dict[str, Any]] = []
    visible_carets: list[dict[str, Any]] = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        if slide_index not in problem_slide_numbers:
            continue
        text_boxes: list[tuple[float, float, float, float]] = []
        rounded: list[tuple[Any, tuple[float, float, float, float]]] = []
        for shape in iter_shapes(slide.shapes):
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                text = shape.text_frame.text
                if "^" in text:
                    visible_carets.append({"slide_number": slide_index, "text": text})
                if text.strip():
                    text_boxes.append(bbox_in(shape))
            try:
                auto_type = str(shape.auto_shape_type).lower()
            except Exception:
                continue
            if "round" in auto_type:
                rounded.append((shape, bbox_in(shape)))

        for shape, bbox in rounded:
            width = shape.width / EMU_PER_INCH
            height = shape.height / EMU_PER_INCH
            area = max(width * height, 0.0001)
            entry = {
                "slide_number": slide_index,
                "x": round(bbox[0], 3),
                "y": round(bbox[1], 3),
                "w": round(width, 3),
                "h": round(height, 3),
            }
            if width <= max_exp_w and height <= max_exp_h:
                small_blanks.append(entry)
                blank_center_y = (bbox[1] + bbox[3]) / 2
                same_row_text = [
                    text_bbox
                    for text_bbox in text_boxes
                    if abs(((text_bbox[1] + text_bbox[3]) / 2) - blank_center_y) < 0.35
                ]
                if same_row_text:
                    nearest_text_top = min(text_bbox[1] for text_bbox in same_row_text)
                    if bbox[1] > nearest_text_top + 0.08:
                        raised_position_issues.append(
                            {
                                **entry,
                                "nearest_text_top": round(nearest_text_top, 3),
                                "top_delta": round(bbox[1] - nearest_text_top, 3),
                            }
                        )
            elif width >= min_base_w and height >= min_base_h:
                baseline_blanks.append(entry)
            for text_bbox in text_boxes:
                ratio = intersection_area(bbox, text_bbox) / area
                if ratio > 0.08:
                    overlaps.append({**entry, "overlap_ratio": round(ratio, 3)})

    problems: list[dict[str, Any]] = []
    if visible_carets:
        problems.append({"code": "visible_caret", "message": "Caret markers are visible in PPTX text.", "details": visible_carets})
    if len(small_blanks) < intent["exponent_blank_count"]:
        problems.append(
            {
                "code": "missing_exponent_blanks",
                "message": "Fewer compact raised blanks were rendered than exponent blank tokens in IR.",
                "expected": intent["exponent_blank_count"],
                "actual": len(small_blanks),
            }
        )
    if len(baseline_blanks) < intent["baseline_blank_count"]:
        problems.append(
            {
                "code": "missing_baseline_blanks",
                "message": "Fewer full-size baseline blanks were rendered than baseline blank tokens in IR.",
                "expected": intent["baseline_blank_count"],
                "actual": len(baseline_blanks),
            }
        )
    if overlaps:
        problems.append({"code": "blank_text_overlap", "message": "Blank shapes overlap nearby text boxes.", "details": overlaps})
    if raised_position_issues:
        problems.append(
            {
                "code": "exponent_blank_not_raised",
                "message": "One or more compact exponent blanks are not raised above the adjacent text row.",
                "details": raised_position_issues,
            }
        )
    if small_blanks and baseline_blanks:
        largest_small = max(item["w"] * item["h"] for item in small_blanks)
        smallest_base = min(item["w"] * item["h"] for item in baseline_blanks)
        if largest_small >= smallest_base * 0.62:
            problems.append(
                {
                    "code": "weak_blank_size_distinction",
                    "message": "Exponent blanks are not visually distinct enough from baseline blanks.",
                    "largest_exponent_area": round(largest_small, 4),
                    "smallest_baseline_area": round(smallest_base, 4),
                }
            )

    return {
        "passes": not problems,
        "intent": intent,
        "small_exponent_blank_shapes": small_blanks,
        "baseline_blank_shapes": baseline_blanks,
        "raised_position_issues": raised_position_issues,
        "problem_count": len(problems),
        "problems": problems,
    }


def next_adjustments(current: dict[str, Any], diagnostics: dict[str, Any]) -> dict[str, Any]:
    adjusted = dict(current)
    codes = {problem["code"] for problem in diagnostics.get("problems", [])}
    scale = float(adjusted.get("math.exponent_blank_scale", 0.48))
    y_offset = float(adjusted.get("math.exponent_blank_y_offset", -0.03))
    gap = float(adjusted.get("math.exponent_blank_gap", 0.04))
    pre_gap = float(adjusted.get("math.exponent_blank_pre_gap", 0.16))
    if "weak_blank_size_distinction" in codes or "blank_text_overlap" in codes:
        adjusted["math.exponent_blank_scale"] = max(0.32, round(scale * 0.88, 3))
        adjusted["math.exponent_blank_pre_gap"] = min(0.16, round(pre_gap + 0.025, 3))
        adjusted["math.exponent_blank_gap"] = min(0.12, round(gap + 0.015, 3))
    if "missing_exponent_blanks" not in codes:
        adjusted["math.exponent_blank_y_offset"] = round(y_offset - 0.015, 3)
    if "exponent_blank_not_raised" in codes:
        adjusted["math.exponent_blank_y_offset"] = round(y_offset - 0.04, 3)
    return adjusted
