from __future__ import annotations

import json
import re
import shutil
import tempfile
import zipfile
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.util import Inches, Pt

from .constants import (
    A_NS,
    BLACK,
    BLANK_LINE,
    COVER_NUMBER_FONT,
    HEADER_BLUE,
    KOREAN_FONT,
    MATH_FONT,
    P_NS,
    SLIDE_H_IN,
    SLIDE_W_IN,
    TITLE_BLUE,
    WHITE,
    YELLOW,
)
from .design import get_design_value, rgb_from_hex
from .ir import PracticeBlock, PracticeItem
from .math_renderer import add_math_row, apply_blank_style
from .ooxml_style import apply_ooxml_style_parts

ET.register_namespace("p", P_NS)
ET.register_namespace("a", A_NS)

BASELINE_BLANK_RE = re.compile(r"(?<!\^)\[\s*\]|__|□")
EMU_PER_INCH = 914400


def style_color(design: dict[str, Any] | None, name: str, default: RGBColor) -> RGBColor:
    return rgb_from_hex(get_design_value(design, f"colors.{name}", None), default)


def style_value(design: dict[str, Any] | None, path: str, default: Any) -> Any:
    return get_design_value(design, path, default)


THEME_COLORS = {
    "BACKGROUND_1": MSO_THEME_COLOR.BACKGROUND_1,
    "BACKGROUND_2": MSO_THEME_COLOR.BACKGROUND_2,
    "TEXT_1": MSO_THEME_COLOR.TEXT_1,
    "TEXT_2": MSO_THEME_COLOR.TEXT_2,
    "ACCENT_1": MSO_THEME_COLOR.ACCENT_1,
    "ACCENT_2": MSO_THEME_COLOR.ACCENT_2,
    "ACCENT_3": MSO_THEME_COLOR.ACCENT_3,
    "ACCENT_4": MSO_THEME_COLOR.ACCENT_4,
    "ACCENT_5": MSO_THEME_COLOR.ACCENT_5,
    "ACCENT_6": MSO_THEME_COLOR.ACCENT_6,
}


def style_theme_color(design: dict[str, Any] | None, path: str) -> MSO_THEME_COLOR | None:
    value = style_value(design, path, None)
    if not value:
        return None
    return THEME_COLORS.get(str(value).strip().upper())


def design_asset_path(design: dict[str, Any] | None, path: str) -> Path | None:
    skill_dir = style_value(design, "_skill_dir", None)
    relative = style_value(design, path, None)
    if not skill_dir or not relative:
        return None
    asset = Path(skill_dir) / relative
    return asset if asset.exists() else None


def add_picture_asset(
    slide,
    asset: Path | None,
    *,
    x: float,
    y: float,
    w: float,
    h: float,
) -> bool:
    if asset is None:
        return False
    slide.shapes.add_picture(str(asset), Inches(x), Inches(y), Inches(w), Inches(h))
    return True


def set_text_frame_style(
    shape,
    font_name: str,
    size: float,
    bold: bool = False,
    color: RGBColor | None = None,
    line_spacing: float = 1.0,
) -> None:
    text_frame = shape.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    text_frame.margin_left = Inches(0.02)
    text_frame.margin_right = Inches(0.02)
    text_frame.margin_top = Inches(0.02)
    text_frame.margin_bottom = Inches(0.02)
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.line_spacing = line_spacing
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(size)
            run.font.bold = bold or font_name.strip().casefold() == MATH_FONT.casefold()
            if color is not None:
                run.font.color.rgb = color


def add_text(
    container,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    *,
    size: float = 24,
    font: str = KOREAN_FONT,
    bold: bool = False,
    color: RGBColor | None = None,
    line_spacing: float = 1.0,
):
    shape = container.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    shape.text = text
    set_text_frame_style(shape, font, size, bold, color or BLACK, line_spacing)
    return shape


def normalize_hex(value: str | None, default: str) -> str:
    if not value:
        return default
    text = value.strip().lstrip("#").upper()
    return text if re.fullmatch(r"[0-9A-F]{6}", text) else default


def text_run_properties(shape) -> list:
    return list(shape._element.findall(f".//{{{A_NS}}}rPr"))


def remove_xml_children(element, tag: str) -> None:
    for child in list(element):
        if child.tag == tag:
            element.remove(child)


def apply_scheme_lum_mod(shape, lum_mod: Any) -> None:
    if lum_mod is None:
        return
    try:
        value = int(lum_mod)
    except (TypeError, ValueError):
        return
    value = max(0, min(100000, value))
    scheme_clr = shape._element.find(f".//{{{A_NS}}}solidFill/{{{A_NS}}}schemeClr")
    if scheme_clr is None:
        return
    remove_xml_children(scheme_clr, f"{{{A_NS}}}lumMod")
    scheme_clr.append(parse_xml(f'<a:lumMod {nsdecls("a")} val="{value}"/>'))


def add_text_outline(shape, *, color: str, width_emu: int) -> None:
    line_tag = f"{{{A_NS}}}ln"
    for r_pr in text_run_properties(shape):
        remove_xml_children(r_pr, line_tag)
        outline = parse_xml(
            f'<a:ln {nsdecls("a")} w="{width_emu}">'
            f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
            f"</a:ln>"
        )
        r_pr.insert(0, outline)


def set_text_fill_scheme(shape, scheme: str) -> None:
    solid_tag = f"{{{A_NS}}}solidFill"
    line_tag = f"{{{A_NS}}}ln"
    for r_pr in text_run_properties(shape):
        remove_xml_children(r_pr, solid_tag)
        fill = parse_xml(f'<a:solidFill {nsdecls("a")}><a:schemeClr val="{scheme}"/></a:solidFill>')
        insert_at = 1 if len(r_pr) and r_pr[0].tag == line_tag else 0
        r_pr.insert(insert_at, fill)


def add_cover_number(slide, number: str, design: dict[str, Any] | None) -> None:
    x = style_value(design, "title_slide.number.x", 2.794)
    y = style_value(design, "title_slide.number.y", 2.527)
    w = style_value(design, "title_slide.number.w", 1.419)
    h = style_value(design, "title_slide.number.h", 1.279)
    size = style_value(design, "title_slide.number.font_size", 70)
    font = style_value(design, "fonts.cover_number", COVER_NUMBER_FONT)
    outline_color = normalize_hex(
        style_value(design, "title_slide.number.outline_color", None),
        "275184",
    )
    outline_width = int(style_value(design, "title_slide.number.outline_width_emu", 152400))
    back_fill = style_color(design, "title_blue", TITLE_BLUE)
    front_scheme = style_value(design, "title_slide.number.front_fill_scheme", "bg1")

    back = add_text(slide.shapes, x, y, w, h, number, size=size, font=font, bold=True, color=back_fill)
    back.text_frame.word_wrap = False
    add_text_outline(back, color=outline_color, width_emu=outline_width)

    front = add_text(slide.shapes, x, y, w, h, number, size=size, font=font, bold=True, color=WHITE)
    front.text_frame.word_wrap = False
    set_text_fill_scheme(front, front_scheme)


def add_practice_header(slide, practice_no: str, design: dict[str, Any] | None) -> None:
    header_asset = design_asset_path(design, f"assets.headers.{practice_no}")
    if header_asset is not None:
        add_picture_asset(
            slide,
            header_asset,
            x=style_value(design, "problem_slide.header.x", 0.28),
            y=style_value(design, "problem_slide.header.y", 0.54),
            w=style_value(design, "problem_slide.header.w", 1.95),
            h=style_value(design, "problem_slide.header.h", 0.46),
        )
        return

    chip = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(style_value(design, "problem_slide.header.x", 0.28)),
        Inches(style_value(design, "problem_slide.header.y", 0.54)),
        Inches(style_value(design, "problem_slide.header.w", 1.95)),
        Inches(style_value(design, "problem_slide.header.h", 0.46)),
    )
    chip.fill.solid()
    chip.fill.fore_color.rgb = style_color(design, "header_blue", HEADER_BLUE)
    chip.line.fill.background()
    chip.shadow.inherit = False
    txt = add_text(
        slide.shapes,
        style_value(design, "problem_slide.header.text_x", 0.38),
        style_value(design, "problem_slide.header.text_y", 0.60),
        style_value(design, "problem_slide.header.text_w", 1.70),
        style_value(design, "problem_slide.header.text_h", 0.28),
        f"개념 익히기 {practice_no}",
        size=style_value(design, "problem_slide.header.font_size", 18),
        bold=True,
        color=WHITE,
    )
    txt.text_frame.word_wrap = False


def add_title_slide(prs: Presentation, concept_no: str, title: str, design: dict[str, Any] | None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if style_value(design, "title_slide.character.enabled", True):
        add_picture_asset(
            slide,
            design_asset_path(design, "assets.title_character"),
            x=style_value(design, "title_slide.character.x", 9.14),
            y=style_value(design, "title_slide.character.y", 2.36),
            w=style_value(design, "title_slide.character.w", 1.53),
            h=style_value(design, "title_slide.character.h", 1.58),
        )
    label = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(style_value(design, "title_slide.label.x", 2.25)),
        Inches(style_value(design, "title_slide.label.y", 2.30)),
        Inches(style_value(design, "title_slide.label.w", 0.78)),
        Inches(style_value(design, "title_slide.label.h", 0.30)),
    )
    label.fill.solid()
    label.fill.fore_color.rgb = style_color(design, "yellow", YELLOW)
    label.line.fill.background()
    add_text(
        slide.shapes,
        style_value(design, "title_slide.label.text_x", 2.31),
        style_value(design, "title_slide.label.text_y", 2.36),
        style_value(design, "title_slide.label.text_w", 0.66),
        style_value(design, "title_slide.label.text_h", 0.20),
        "개념핵심",
        size=style_value(design, "title_slide.label.font_size", 10.5),
        bold=True,
        color=style_color(design, "title_blue", TITLE_BLUE),
    )
    number = concept_no or "00"
    add_cover_number(slide, number, design)
    add_text(
        slide.shapes,
        style_value(design, "title_slide.title.x", 4.05),
        style_value(design, "title_slide.title.y", 2.82),
        style_value(design, "title_slide.title.w", 7.20),
        style_value(design, "title_slide.title.h", 0.60),
        title,
        size=style_value(design, "title_slide.title.font_size", 38),
        bold=True,
    )


def split_items(items: list[PracticeItem], chunk_size: int = 4) -> list[list[PracticeItem]]:
    return [items[idx : idx + chunk_size] for idx in range(0, len(items), chunk_size)]


def item_display_line_count(item: PracticeItem) -> int:
    return max(1, len([line for line in item.display_lines if line.strip()]))


def split_multiline_vertical_items(
    items: list[PracticeItem],
    line_budget: int = 7,
    max_items_per_slide: int = 4,
) -> list[list[PracticeItem]]:
    chunks: list[list[PracticeItem]] = []
    current: list[PracticeItem] = []
    current_cost = 0
    for item in items:
        cost = item_display_line_count(item)
        force_single_item_slide = cost >= 5 and item.blank_count >= 4
        if current and force_single_item_slide:
            chunks.append(current)
            current = []
            current_cost = 0
        if current and (current_cost + cost > line_budget or len(current) >= max_items_per_slide):
            chunks.append(current)
            current = []
            current_cost = 0
        current.append(item)
        current_cost += cost
        if force_single_item_slide or current_cost >= line_budget:
            chunks.append(current)
            current = []
            current_cost = 0
    if current:
        chunks.append(current)
    return chunks


def split_block_items(block: PracticeBlock) -> list[list[PracticeItem]]:
    items = block.item_models
    if block.page == 10 and block.practice_no == "1":
        return [items[:3], items[3:]]
    if block.page == 11 and block.practice_no == "2":
        return [items[:2], items[2:]]
    if block.page == 14 and block.practice_no == "1":
        return [items[:2], items[2:]]
    if block.layout_type == "two_column_grid" and any(is_multiline_item(item) for item in items):
        return split_two_column_worked_rows(items)
    if block.layout_type != "two_column_grid" and any(is_multiline_item(item) for item in items):
        return split_multiline_vertical_items(items)
    return split_items(items)


def display_order(items: list[PracticeItem], layout_type: str) -> list[PracticeItem]:
    if layout_type == "two_column_grid":
        return sorted(
            items,
            key=lambda item: (
                item.row_index,
                item.column_index,
                item.number if item.number is not None else 999,
            ),
        )
    return items


def is_multiline_item(item: PracticeItem) -> bool:
    display_line_count = len([line for line in item.display_lines if line.strip()])
    if display_line_count >= 2:
        return True
    segment_line_indexes = {
        int(getattr(segment, "line_index", 0))
        for segment in item.display_segments
        if getattr(segment, "text", "").strip()
    }
    if item.display_segments and len(segment_line_indexes) <= 1:
        return False
    expr = item.expression_text
    source_line_count = len([line for line in item.source_lines if line.strip()])
    if source_line_count >= 3:
        return True
    return source_line_count >= 2 and len(expr) >= 68


def normalize_spacing_preserving_parenthesis_blanks(text: str) -> str:
    placeholders: list[str] = []

    def stash(match: re.Match[str]) -> str:
        placeholders.append(match.group(0))
        return f"\uE000{len(placeholders) - 1}\uE000"

    normalized = re.sub(r"\(\s+\)", stash, text)
    normalized = re.sub(r"\s+", " ", normalized).strip()
    for idx, value in enumerate(placeholders):
        normalized = normalized.replace(f"\uE000{idx}\uE000", value)
    return normalized


def split_worked_expression_lines(item: PracticeItem) -> list[str]:
    display_lines = [line.strip() for line in item.display_lines if line.strip()]
    if len(display_lines) >= 2:
        return display_lines

    expr = item.expression_text.strip()
    if not is_multiline_item(item):
        return [expr]

    expr = normalize_spacing_preserving_parenthesis_blanks(expr)
    parts = re.split(r"\s*=\s*", expr)
    if len(parts) <= 1:
        return [expr]

    lines = [parts[0].strip()]
    for part in parts[1:]:
        part = part.strip()
        if part:
            lines.append("=" + part)
    return lines


def display_segment_rows(item: PracticeItem) -> list[list[Any]]:
    rows: dict[int, list[Any]] = {}
    for segment in item.display_segments:
        text = getattr(segment, "text", "").strip()
        kind = getattr(segment, "kind", "math")
        if not text and kind not in {"marker", "number_line", "spacer", "blank_shape"}:
            continue
        line_index = int(getattr(segment, "line_index", 0))
        rows.setdefault(line_index, []).append(segment)
    return [rows[line_index] for line_index in sorted(rows)]


def display_segments_are_structured(item: PracticeItem, rows: list[list[Any]]) -> bool:
    if len([line for line in item.display_lines if line.strip()]) >= 2:
        return True
    if any(getattr(segment, "gap_after_in", None) is not None for row in rows for segment in row):
        return True
    if any(
        getattr(segment, "kind", "math") in {"marker", "number_line", "value_table", "equation_system", "blank_shape"}
        for row in rows
        for segment in row
    ):
        return True
    return any(getattr(segment, "kind", "math") == "korean_label" for row in rows for segment in row)


def item_display_texts(item: PracticeItem) -> list[str]:
    texts = [line for line in item.display_lines if line.strip()]
    texts.extend(
        getattr(segment, "text", "")
        for segment in item.display_segments
        if getattr(segment, "text", "").strip()
    )
    return texts or [item.expression_text, item.raw_text]


def worked_item_line_gap(item: PracticeItem, design: dict[str, Any] | None, base_gap: float) -> float:
    texts = item_display_texts(item)
    has_baseline_blank = any(BASELINE_BLANK_RE.search(text) for text in texts)
    if not has_baseline_blank:
        return base_gap

    has_fraction_blank = any("/" in text and BASELINE_BLANK_RE.search(text) for text in texts)
    if has_fraction_blank:
        return max(
            base_gap,
            style_value(design, "problem_slide.worked_item.fraction_blank_line_gap", 0.92),
        )

    blank_h = style_value(design, "problem_slide.item.blank_h", 0.551)
    blank_y_offset = style_value(design, "problem_slide.item.blank_y_offset", 0.03)
    blank_margin = style_value(design, "problem_slide.worked_item.blank_line_margin", 0.14)
    minimum_gap = blank_h + blank_y_offset + blank_margin
    return max(
        base_gap,
        style_value(design, "problem_slide.worked_item.blank_line_gap", minimum_gap),
    )


def segment_gap_after(segment: Any, default_gap: float) -> float:
    value = getattr(segment, "gap_after_in", None)
    if value is None:
        return default_gap
    try:
        return max(0.0, float(value))
    except (TypeError, ValueError):
        return default_gap


def segment_width(segment: Any, default_width: float) -> float:
    value = getattr(segment, "width_in", None)
    if value is None:
        return default_width
    try:
        return max(0.0, float(value))
    except (TypeError, ValueError):
        return default_width


def add_korean_label_text(
    container,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    *,
    design: dict[str, Any] | None,
    size: float,
) -> None:
    label = add_text(
        container,
        x,
        y,
        w,
        h,
        text,
        size=size,
        font=style_value(design, "fonts.korean", KOREAN_FONT),
        bold=True,
    )
    label.text_frame.word_wrap = False


def add_marker_shape(
    container,
    x: float,
    y: float,
    segment: Any,
    *,
    design: dict[str, Any] | None,
) -> float:
    shape_name = str(getattr(segment, "shape", "") or "right_arrow").strip()
    if shape_name != "right_arrow":
        return 0.0

    marker_w = style_value(design, "problem_slide.item.marker_right_arrow_w", 0.4)
    marker_h = style_value(design, "problem_slide.item.marker_right_arrow_h", 0.292)
    marker_y = y + style_value(design, "problem_slide.item.marker_right_arrow_y_offset", 0.101)
    marker = container.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(x),
        Inches(marker_y),
        Inches(marker_w),
        Inches(marker_h),
    )
    marker.fill.solid()
    fill_theme = style_theme_color(design, "problem_slide.item.marker_right_arrow_fill_theme_color")
    if fill_theme is not None:
        marker.fill.fore_color.theme_color = fill_theme
        apply_scheme_lum_mod(
            marker,
            style_value(design, "problem_slide.item.marker_right_arrow_fill_lum_mod", 50000),
        )
    else:
        marker.fill.fore_color.rgb = style_color(design, "blank_line", BLANK_LINE)
    marker.line.fill.background()
    return marker_w


def parse_number_line_spec(text: str) -> dict[str, str]:
    spec: dict[str, str] = {}
    for part in text.split(";"):
        if "=" not in part:
            continue
        key, value = part.split("=", 1)
        spec[key.strip()] = value.strip()
    return spec


def parse_float_value(value: str | None, default: float) -> float:
    if value is None or value == "":
        return default
    try:
        return float(value)
    except ValueError:
        return default


def parse_bool_value(value: str | None, default: bool = False) -> bool:
    if value is None:
        return default
    return value.strip().lower() in {"1", "true", "yes", "y", "closed", "solid"}


def parse_number_line_values(value: str | None) -> list[float]:
    if not value:
        return []
    values: list[float] = []
    for token in value.split(","):
        token = token.strip()
        if not token:
            continue
        try:
            values.append(float(token))
        except ValueError:
            continue
    return values


def add_solid_rect(
    container,
    x: float,
    y: float,
    w: float,
    h: float,
    color: RGBColor,
    alpha: float | None = None,
):
    rect = container.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    if alpha is not None:
        rect.fill.transparency = max(0.0, min(1.0, alpha))
    rect.line.fill.background()
    return rect


def add_number_line_arrowhead(container, x: float, y: float, direction: str, color: RGBColor) -> None:
    arrow_w = 0.105
    arrow_h = 0.105
    triangle = container.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        Inches(x),
        Inches(y),
        Inches(arrow_w),
        Inches(arrow_h),
    )
    triangle.rotation = 90 if direction == "right" else 270
    triangle.fill.solid()
    triangle.fill.fore_color.rgb = color
    triangle.line.fill.background()


def add_number_line_segment(
    container,
    x: float,
    y: float,
    w: float,
    segment: Any,
    *,
    design: dict[str, Any] | None,
) -> float:
    spec = parse_number_line_spec(getattr(segment, "text", ""))
    point = parse_float_value(spec.get("point"), 0.0)
    min_value = parse_float_value(spec.get("min"), point - 2.0)
    max_value = parse_float_value(spec.get("max"), point + 2.0)
    if min_value == max_value:
        max_value = min_value + 1.0

    axis_color = style_color(design, "text", BLACK)
    shade_color = RGBColor(190, 153, 203)
    axis_y = y + 0.31
    axis_h = 0.018
    arrow_w = 0.105
    tick_h = 0.13
    tick_w = 0.012
    line_x0 = x + arrow_w
    line_x1 = x + w - arrow_w
    usable_w = max(0.1, line_x1 - line_x0)

    def pos(value: float) -> float:
        return line_x0 + ((value - min_value) / (max_value - min_value)) * usable_w

    add_solid_rect(container, line_x0, axis_y, usable_w, axis_h, axis_color)
    add_number_line_arrowhead(container, x, axis_y - 0.044, "left", axis_color)
    add_number_line_arrowhead(container, x + w - arrow_w, axis_y - 0.044, "right", axis_color)

    ticks = parse_number_line_values(spec.get("ticks"))
    if not ticks and spec.get("blank", "").lower() == "true":
        ticks = [float(value) for value in range(int(min_value), int(max_value) + 1)]
    label_values = parse_number_line_values(spec.get("labels")) or ticks
    labels = {
        round(value, 6): str(int(value)) if float(value).is_integer() else str(value)
        for value in label_values
    }

    for tick in ticks:
        tick_x = pos(tick)
        add_solid_rect(container, tick_x - tick_w / 2, axis_y - tick_h / 2, tick_w, tick_h, axis_color)
        label = labels.get(round(tick, 6))
        if label:
            add_text(
                container,
                tick_x - 0.22,
                axis_y + 0.08,
                0.44,
                0.20,
                label,
                size=style_value(design, "problem_slide.number_line.label_font_size", 10.5),
                font=style_value(design, "fonts.math", MATH_FONT),
                bold=True,
            )

    direction = spec.get("direction", "").strip().lower()
    if direction in {"left", "right"} and "point" in spec:
        point_x = pos(point)
        shade_y = axis_y - 0.22
        shade_h = 0.20
        shade_x0, shade_x1 = (point_x, line_x1) if direction == "right" else (line_x0, point_x)
        if shade_x1 > shade_x0:
            add_solid_rect(container, shade_x0, shade_y, shade_x1 - shade_x0, shade_h, shade_color, alpha=0.25)
            add_solid_rect(container, shade_x0, shade_y, shade_x1 - shade_x0, 0.024, axis_color)
            arrow_x = shade_x1 - arrow_w if direction == "right" else shade_x0
            add_number_line_arrowhead(container, arrow_x, shade_y - 0.041, direction, axis_color)
            add_solid_rect(container, point_x - 0.011, shade_y, 0.022, axis_y - shade_y, axis_color)

        dot_d = 0.14
        dot = container.add_shape(
            MSO_SHAPE.OVAL,
            Inches(point_x - dot_d / 2),
            Inches(axis_y + axis_h / 2 - dot_d / 2),
            Inches(dot_d),
            Inches(dot_d),
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = axis_color if parse_bool_value(spec.get("closed")) else WHITE
        dot.line.color.rgb = axis_color
        dot.line.width = Pt(1.0)

    return w


def segment_spec(segment: Any) -> dict[str, Any]:
    text = str(getattr(segment, "text", "") or "").strip()
    if not text:
        return {}
    try:
        value = json.loads(text)
    except json.JSONDecodeError:
        return {}
    return value if isinstance(value, dict) else {}


def add_centered_text(
    container,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    *,
    size: float,
    font: str,
    bold: bool = True,
    color: RGBColor | None = None,
) -> None:
    shape = add_text(container, x, y, w, h, text, size=size, font=font, bold=bold, color=color or BLACK)
    shape.text_frame.word_wrap = False
    for paragraph in shape.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER


def simple_fraction_parts(text: str) -> tuple[str, str] | None:
    match = re.fullmatch(r"\s*(-?\d+)\s*/\s*(-?\d+)\s*", text)
    if not match:
        return None
    return match.group(1), match.group(2)


def add_table_cell_fraction(
    container,
    x: float,
    y: float,
    w: float,
    h: float,
    numerator: str,
    denominator: str,
    *,
    size: float,
    font: str,
    color: RGBColor | None = None,
) -> None:
    fraction_w = min(max(0.16, w * 0.34), max(0.12, w - 0.08))
    fraction_x = x + (w - fraction_w) / 2
    bar_y = y + h * 0.50
    num_y = y + h * 0.03
    den_y = y + h * 0.50
    text_h = max(0.08, h * 0.42)
    fraction_size = max(6.5, size * 0.82)
    add_centered_text(
        container,
        fraction_x,
        num_y,
        fraction_w,
        text_h,
        numerator,
        size=fraction_size,
        font=font,
        color=color,
    )
    add_solid_rect(container, fraction_x + 0.01, bar_y, max(0.04, fraction_w - 0.02), 0.006, color or BLACK)
    add_centered_text(
        container,
        fraction_x,
        den_y,
        fraction_w,
        text_h,
        denominator,
        size=fraction_size,
        font=font,
        color=color,
    )


def add_table_cell_text(
    container,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    *,
    size: float,
    font: str,
    color: RGBColor | None = None,
) -> None:
    fraction = simple_fraction_parts(text)
    if fraction:
        add_table_cell_fraction(container, x, y, w, h, fraction[0], fraction[1], size=size, font=font, color=color)
        return
    add_centered_text(container, x, y + 0.02, w, max(0.1, h - 0.03), text, size=size, font=font, color=color)


def add_value_table_segment(
    container,
    x: float,
    y: float,
    segment: Any,
    *,
    design: dict[str, Any] | None,
    size: float,
) -> float:
    spec = segment_spec(segment)
    row_labels = [str(value) for value in spec.get("row_labels", ["x", "y"])]
    rows = spec.get("rows", [])
    if not isinstance(rows, list) or not rows:
        return 0.0
    rows = [[str(cell) for cell in row] if isinstance(row, list) else [] for row in rows]
    column_count = max((len(row) for row in rows), default=0)
    if not column_count:
        return 0.0

    row_label_w = float(spec.get("row_label_w", 0.36) or 0.36)
    cell_w = float(spec.get("cell_w", 0.46) or 0.46)
    cell_h = float(spec.get("cell_h", 0.29) or 0.29)
    table_font_size = float(spec.get("font_size", max(8.5, size * 0.62)) or max(8.5, size * 0.62))
    line_color = rgb_from_hex(str(spec.get("line_color", "#888888")), RGBColor(136, 136, 136))
    header_fill = rgb_from_hex(str(spec.get("header_fill", "#DDEED7")), RGBColor(221, 238, 215))
    white_fill = RGBColor(255, 255, 255)
    red_cells = {
        (int(cell[0]), int(cell[1]))
        for cell in spec.get("red_cells", [])
        if isinstance(cell, list) and len(cell) == 2
    }
    hide_red_cells = bool(spec.get("hide_red_cells", False))
    math_font = style_value(design, "fonts.math", MATH_FONT)

    table_group = container.add_group_shape()
    table_group.name = str(spec.get("name", "value-table"))

    for row_idx, row in enumerate(rows):
        row_y = y + row_idx * cell_h
        label_text = row_labels[row_idx] if row_idx < len(row_labels) else ""
        for col_idx in range(column_count + 1):
            cell_x = x if col_idx == 0 else x + row_label_w + (col_idx - 1) * cell_w
            cell_w_current = row_label_w if col_idx == 0 else cell_w
            rect = table_group.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(cell_x),
                Inches(row_y),
                Inches(cell_w_current),
                Inches(cell_h),
            )
            rect.fill.solid()
            rect.fill.fore_color.rgb = header_fill if col_idx == 0 else white_fill
            rect.line.color.rgb = line_color
            rect.line.width = Pt(0.5)

            text = label_text if col_idx == 0 else (row[col_idx - 1] if col_idx - 1 < len(row) else "")
            if hide_red_cells and (row_idx, col_idx - 1) in red_cells:
                text = ""
            if text:
                add_table_cell_text(
                    table_group.shapes,
                    cell_x,
                    row_y,
                    cell_w_current,
                    cell_h,
                    text,
                    size=table_font_size,
                    font=math_font,
                    color=RGBColor(255, 0, 0) if (row_idx, col_idx - 1) in red_cells and not hide_red_cells else BLACK,
                )

    return row_label_w + column_count * cell_w


def add_equation_system_segment(
    container,
    x: float,
    y: float,
    segment: Any,
    *,
    design: dict[str, Any] | None,
    size: float,
) -> float:
    spec = segment_spec(segment)
    equations = [str(value) for value in spec.get("equations", []) if str(value).strip()]
    if not equations:
        return 0.0
    suffix = str(spec.get("suffix", "") or "")
    line_gap = float(spec.get("line_gap", 0.25) or 0.25)
    brace_w = float(spec.get("brace_w", 0.20) or 0.20)
    equation_w = float(spec.get("equation_w", 1.05) or 1.05)
    system_size = float(spec.get("font_size", max(10.5, size * 0.68)) or max(10.5, size * 0.68))
    math_font = style_value(design, "fonts.math", MATH_FONT)
    korean_font = style_value(design, "fonts.korean", KOREAN_FONT)
    height = max(0.34, line_gap * max(1, len(equations)))

    system_group = container.add_group_shape()
    system_group.name = str(spec.get("name", "equation-system"))
    add_text(
        system_group.shapes,
        x,
        y - 0.03,
        brace_w,
        height + 0.12,
        "{",
        size=system_size * 1.9,
        font=math_font,
        bold=True,
    )
    for idx, equation in enumerate(equations):
        add_math_row(
            system_group.shapes,
            x + brace_w * 0.72,
            y + idx * line_gap,
            equation_w,
            0.23,
            equation,
            design=design,
            size=system_size,
            font=math_font,
        )
    if suffix:
        add_text(
            system_group.shapes,
            x + brace_w * 0.72 + equation_w + 0.06,
            y + max(0.0, (height - 0.22) / 2),
            float(spec.get("suffix_w", 0.72) or 0.72),
            0.24,
            suffix,
            size=max(8.5, system_size * 0.78),
            font=korean_font,
            bold=True,
        )
    return brace_w * 0.72 + equation_w + (float(spec.get("suffix_w", 0.72) or 0.72) if suffix else 0.0) + 0.08


def add_blank_shape_segment(
    container,
    x: float,
    y: float,
    segment: Any,
    *,
    design: dict[str, Any] | None,
) -> float:
    spec = segment_spec(segment)
    scale = float(spec.get("scale", 0.55) or 0.55)
    blank_w = float(spec.get("w", style_value(design, "problem_slide.item.blank_w", 0.551) * scale))
    blank_h = float(spec.get("h", style_value(design, "problem_slide.item.blank_h", 0.551) * scale))
    blank_y_offset = float(spec.get("y_offset", 0.05) or 0.05)
    blank = container.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y + blank_y_offset), Inches(blank_w), Inches(blank_h))
    blank.name = str(spec.get("name", "answer-blank"))
    apply_blank_style(blank, design)
    return blank_w


def add_display_segment_row(
    container,
    x: float,
    y: float,
    w: float,
    h: float,
    segments: list[Any],
    *,
    design: dict[str, Any] | None,
    size: float,
) -> None:
    cursor = x
    label_gap = style_value(design, "problem_slide.item.inline_label_gap", 0.28)
    label_w = style_value(design, "problem_slide.item.inline_label_w", 2.25)
    for idx, segment in enumerate(segments):
        kind = getattr(segment, "kind", "math")
        if kind == "spacer":
            cursor += segment_width(segment, 0.0) + segment_gap_after(segment, 0.0)
            continue
        if kind == "value_table":
            used_w = add_value_table_segment(container, cursor, y, segment, design=design, size=size)
            cursor += used_w + segment_gap_after(segment, label_gap)
            continue
        if kind == "equation_system":
            used_w = add_equation_system_segment(container, cursor, y, segment, design=design, size=size)
            cursor += used_w + segment_gap_after(segment, label_gap)
            continue
        if kind == "blank_shape":
            used_w = add_blank_shape_segment(container, cursor, y, segment, design=design)
            cursor += used_w + segment_gap_after(segment, 0.08)
            continue
        if kind == "number_line":
            used_w = add_number_line_segment(container, cursor, y, min(3.4, x + w - cursor), segment, design=design)
            cursor += used_w + segment_gap_after(segment, label_gap)
            continue
        if kind == "marker":
            used_w = add_marker_shape(container, cursor, y, segment, design=design)
            cursor += used_w + segment_gap_after(
                segment,
                style_value(design, "problem_slide.item.marker_right_arrow_gap_after", label_gap),
            )
            continue

        text = getattr(segment, "text", "").strip()
        if not text or cursor >= x + w:
            continue
        if kind == "korean_label":
            width = min(segment_width(segment, label_w), max(0.1, x + w - cursor))
            add_korean_label_text(container, cursor, y, width, h, text, design=design, size=size)
            cursor += width + segment_gap_after(segment, label_gap)
            continue

        later_label_count = sum(
            1 for later in segments[idx + 1 :] if getattr(later, "kind", "math") == "korean_label"
        )
        reserved_width = later_label_count * (label_w + label_gap)
        math_w = max(0.1, x + w - cursor - reserved_width)
        used_w = add_math_row(
            container,
            cursor,
            y,
            math_w,
            h,
            text,
            design=design,
            size=size,
            font=style_value(design, "fonts.math", MATH_FONT),
        )
        cursor += used_w + segment_gap_after(segment, label_gap)


def layout_shape_x(
    shape: Any,
    *,
    number_x: float,
    formula_x: float,
) -> float:
    anchor = str(getattr(shape, "x_anchor", "") or "before_display_lines").strip()
    offset = float(getattr(shape, "x_offset_in", 0.0) or 0.0)
    if anchor == "before_number":
        return number_x + offset
    if anchor == "absolute":
        return offset
    return formula_x + offset


def line_shape_color(shape: Any, design: dict[str, Any] | None) -> RGBColor:
    return rgb_from_hex(
        str(getattr(shape, "stroke_color", "") or "").strip(),
        style_color(design, "text", BLACK),
    )


OOXML_DASH_VALUES = {
    "solid": "solid",
    "dash": "dash",
    "long_dash": "lgDash",
    "round_dot": "sysDot",
    "square_dot": "sysDash",
    "dash_dot": "dashDot",
}


def inches_to_emu(value: float) -> int:
    return int(round(float(value) * EMU_PER_INCH))


def point_or_default(value: Any, default: float) -> float:
    try:
        return float(value)
    except (TypeError, ValueError):
        return default


def curve_points(shape: Any) -> tuple[tuple[float, float], tuple[float, float], tuple[float, float], tuple[float, float]] | None:
    x1 = getattr(shape, "x1_in", None)
    y1 = getattr(shape, "y1_in", None)
    x2 = getattr(shape, "x2_in", None)
    y2 = getattr(shape, "y2_in", None)
    if None in {x1, y1, x2, y2}:
        return None
    x1 = float(x1)
    y1 = float(y1)
    x2 = float(x2)
    y2 = float(y2)
    c1x = point_or_default(getattr(shape, "control1_x_in", None), x1 + (x2 - x1) / 3)
    c1y = point_or_default(getattr(shape, "control1_y_in", None), y1 + (y2 - y1) / 3)
    c2x = point_or_default(getattr(shape, "control2_x_in", None), x1 + 2 * (x2 - x1) / 3)
    c2y = point_or_default(getattr(shape, "control2_y_in", None), y1 + 2 * (y2 - y1) / 3)
    return (x1, y1), (c1x, c1y), (c2x, c2y), (x2, y2)


def set_freeform_cubic_path(sp, points: tuple[tuple[float, float], tuple[float, float], tuple[float, float], tuple[float, float]]) -> None:
    xs = [point[0] for point in points]
    ys = [point[1] for point in points]
    min_x = min(xs)
    min_y = min(ys)
    max_x = max(xs)
    max_y = max(ys)
    width = max(0.01, max_x - min_x)
    height = max(0.01, max_y - min_y)

    xfrm = sp.find(f"./{{{P_NS}}}spPr/{{{A_NS}}}xfrm")
    off = xfrm.find(f"./{{{A_NS}}}off") if xfrm is not None else None
    ext = xfrm.find(f"./{{{A_NS}}}ext") if xfrm is not None else None
    if off is not None:
        off.set("x", str(inches_to_emu(min_x)))
        off.set("y", str(inches_to_emu(min_y)))
    if ext is not None:
        ext.set("cx", str(inches_to_emu(width)))
        ext.set("cy", str(inches_to_emu(height)))

    rel_points = [
        (inches_to_emu(x - min_x), inches_to_emu(y - min_y))
        for x, y in points
    ]
    path_w = inches_to_emu(width)
    path_h = inches_to_emu(height)
    path_xml = (
        f'<a:pathLst {nsdecls("a")}>'
        f'<a:path w="{path_w}" h="{path_h}" fill="none" stroke="true">'
        f'<a:moveTo><a:pt x="{rel_points[0][0]}" y="{rel_points[0][1]}"/></a:moveTo>'
        "<a:cubicBezTo>"
        f'<a:pt x="{rel_points[1][0]}" y="{rel_points[1][1]}"/>'
        f'<a:pt x="{rel_points[2][0]}" y="{rel_points[2][1]}"/>'
        f'<a:pt x="{rel_points[3][0]}" y="{rel_points[3][1]}"/>'
        "</a:cubicBezTo>"
        "</a:path>"
        "</a:pathLst>"
    )
    cust_geom = sp.find(f"./{{{P_NS}}}spPr/{{{A_NS}}}custGeom")
    if cust_geom is None:
        return
    existing_path_lst = cust_geom.find(f"./{{{A_NS}}}pathLst")
    if existing_path_lst is not None:
        cust_geom.remove(existing_path_lst)
    cust_geom.append(parse_xml(path_xml))


def set_freeform_line_style(sp, shape: Any, design: dict[str, Any] | None) -> None:
    sp_pr = sp.find(f"./{{{P_NS}}}spPr")
    if sp_pr is None:
        return
    for child in list(sp_pr):
        if child.tag in {f"{{{A_NS}}}noFill", f"{{{A_NS}}}solidFill", f"{{{A_NS}}}ln"}:
            sp_pr.remove(child)
    sp_pr.append(parse_xml(f'<a:noFill {nsdecls("a")}/>'))
    stroke = line_shape_color(shape, design)
    color_hex = f"{stroke[0]:02X}{stroke[1]:02X}{stroke[2]:02X}"
    width_emu = max(1, int(round(float(getattr(shape, "stroke_pt", 1.0) or 1.0) * 12700)))
    dash = OOXML_DASH_VALUES.get(str(getattr(shape, "stroke_dash", "solid") or "solid").strip().lower(), "solid")
    dash_xml = "" if dash == "solid" else f'<a:prstDash val="{dash}"/>'
    arrowhead = str(getattr(shape, "arrowhead", "") or "").strip().lower()
    head_xml = ""
    if arrowhead not in {"", "none", "false"}:
        head_type = "triangle" if arrowhead in {"triangle", "arrow", "arrowhead"} else arrowhead
        head_xml = f'<a:headEnd type="{head_type}" w="sm" len="sm"/>'
    sp_pr.append(
        parse_xml(
            f'<a:ln {nsdecls("a")} w="{width_emu}">'
            f'<a:solidFill><a:srgbClr val="{color_hex}"/></a:solidFill>'
            f"{dash_xml}"
            f"{head_xml}"
            "</a:ln>"
        )
    )


def add_curved_arrow_shape(
    container,
    shape: Any,
    *,
    number_x: float,
    formula_x: float,
    y: float,
    line_gap: float,
    formula_y_offset: float,
    design: dict[str, Any] | None,
) -> bool:
    points = curve_points(shape)
    if points is None:
        start_x = layout_shape_x(shape, number_x=number_x, formula_x=formula_x)
        start_y = y + formula_y_offset + float(getattr(shape, "y_start_offset_in", 0.0) or 0.0)
        end_x = start_x + float(getattr(shape, "width_in", 0.16) or 0.16)
        end_y = y + formula_y_offset + line_gap + float(getattr(shape, "y_end_offset_in", 0.0) or 0.0)
        points = (
            (start_x, start_y),
            (start_x + (end_x - start_x) / 3, start_y + (end_y - start_y) / 3),
            (start_x + 2 * (end_x - start_x) / 3, start_y + 2 * (end_y - start_y) / 3),
            (end_x, end_y),
        )

    xs = [point[0] for point in points]
    ys = [point[1] for point in points]
    sp = container._spTree.add_freeform_sp(
        inches_to_emu(min(xs)),
        inches_to_emu(min(ys)),
        inches_to_emu(max(0.01, max(xs) - min(xs))),
        inches_to_emu(max(0.01, max(ys) - min(ys))),
    )
    c_nv_pr = sp.find(f"./{{{P_NS}}}nvSpPr/{{{P_NS}}}cNvPr")
    if c_nv_pr is not None:
        c_nv_pr.set("name", str(getattr(shape, "name", "") or "curved-arrow"))
    set_freeform_cubic_path(sp, points)
    set_freeform_line_style(sp, shape, design)
    return True


def add_brace_connector_shape(
    container,
    shape: Any,
    *,
    number_x: float,
    formula_x: float,
    y: float,
    line_gap: float,
    formula_y_offset: float,
    design: dict[str, Any] | None,
) -> bool:
    shape_name = str(getattr(shape, "shape", "") or "left_square_bracket").strip()
    if shape_name not in {"left_square_bracket", "right_square_bracket"}:
        return False

    try:
        line_start = int(getattr(shape, "line_start", 0))
    except (TypeError, ValueError):
        line_start = 0
    try:
        line_end = int(getattr(shape, "line_end", line_start))
    except (TypeError, ValueError):
        line_end = line_start
    if line_end < line_start:
        line_start, line_end = line_end, line_start

    tick_lines = []
    for line in getattr(shape, "tick_lines", []) or []:
        try:
            line_value = int(line)
        except (TypeError, ValueError):
            continue
        if line_start <= line_value <= line_end:
            tick_lines.append(line_value)
    if not tick_lines:
        tick_lines = sorted({line_start, line_end})

    width = max(0.02, float(getattr(shape, "width_in", 0.16) or 0.16))
    stroke_w = max(0.003, float(getattr(shape, "stroke_pt", 1.0) or 1.0) / 72.0)
    color = line_shape_color(shape, design)
    x0 = layout_shape_x(shape, number_x=number_x, formula_x=formula_x)
    vertical_x = x0 if shape_name == "left_square_bracket" else x0 + width - stroke_w
    tick_x = x0 if shape_name == "left_square_bracket" else x0

    y_start = (
        y
        + formula_y_offset
        + line_start * line_gap
        + float(getattr(shape, "y_start_offset_in", 0.0) or 0.0)
    )
    y_end = (
        y
        + formula_y_offset
        + line_end * line_gap
        + float(getattr(shape, "y_end_offset_in", 0.0) or 0.0)
    )
    if y_end < y_start:
        y_start, y_end = y_end, y_start
    add_solid_rect(container, vertical_x, y_start, stroke_w, max(stroke_w, y_end - y_start + stroke_w), color)

    tick_y_offset = float(getattr(shape, "tick_y_offset_in", 0.0) or 0.0)
    for line in tick_lines:
        tick_y = y + formula_y_offset + line * line_gap + tick_y_offset
        add_solid_rect(container, tick_x, tick_y, width, stroke_w, color)
    return True


def add_item_layout_shapes(
    container,
    item: PracticeItem,
    *,
    number_x: float,
    formula_x: float,
    y: float,
    line_gap: float,
    formula_y_offset: float,
    design: dict[str, Any] | None,
) -> int:
    rendered_count = 0
    for shape in item.layout_shapes:
        kind = getattr(shape, "kind", "")
        if kind == "brace_connector":
            if add_brace_connector_shape(
                container,
                shape,
                number_x=number_x,
                formula_x=formula_x,
                y=y,
                line_gap=line_gap,
                formula_y_offset=formula_y_offset,
                design=design,
            ):
                rendered_count += 1
        elif kind == "curved_arrow":
            if add_curved_arrow_shape(
                container,
                shape,
                number_x=number_x,
                formula_x=formula_x,
                y=y,
                line_gap=line_gap,
                formula_y_offset=formula_y_offset,
                design=design,
            ):
                rendered_count += 1
    return rendered_count


def split_two_column_worked_rows(items: list[PracticeItem]) -> list[list[PracticeItem]]:
    ordered_items = display_order(items, "two_column_grid")
    if not ordered_items:
        return []
    rows: dict[int, list[PracticeItem]] = {}
    for fallback_idx, item in enumerate(ordered_items):
        row = item.row_index if item.source_box else fallback_idx // 2
        rows.setdefault(row, []).append(item)
    return [rows[row] for row in sorted(rows)]


def add_blank_if_needed(group, text: str, x: float, y: float, design: dict[str, Any] | None) -> None:
    blank_count = len(re.findall(r"\^\[\s*\]|\[\s*\]|__|□", text))
    if blank_count == 0 and "빈칸" in text:
        blank_count = 1
    for _ in range(blank_count):
        blank = group.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(style_value(design, "problem_slide.item.blank_w", 0.55)),
            Inches(style_value(design, "problem_slide.item.blank_h", 0.42)),
        )
        apply_blank_style(blank, design)
        x += style_value(design, "problem_slide.item.blank_gap", 0.68)


def item_group_label(item: PracticeItem, idx: int) -> str:
    number = item.number if item.number is not None else idx + 1
    return f"problem-{number}"


def add_item_group(
    slide,
    item: PracticeItem,
    idx: int,
    *,
    reveal: bool,
    design: dict[str, Any] | None,
    number_x: float | None = None,
    formula_x: float | None = None,
    y: float | None = None,
    formula_w: float | None = None,
    formula_h: float | None = None,
) -> str | None:
    item_text = item.raw_text
    if y is None:
        y = style_value(design, "problem_slide.item.start_y", 2.05) + idx * style_value(design, "problem_slide.item.row_gap", 1.12)
    if number_x is None:
        number_x = style_value(design, "problem_slide.item.number_x", 0.40)
    if formula_x is None:
        formula_x = style_value(design, "problem_slide.item.formula_x", 1.05)
    if formula_w is None:
        formula_w = style_value(design, "problem_slide.item.formula_w", 7.20)
    if formula_h is None:
        formula_h = style_value(design, "problem_slide.item.formula_h", 0.58)
    multiline = is_multiline_item(item)
    display_line_count = item_display_line_count(item)
    formula_size = style_value(design, "problem_slide.item.formula_font_size", 21)
    if multiline:
        formula_size = style_value(design, "problem_slide.worked_item.formula_font_size", min(18, formula_size))
        formula_h = style_value(design, "problem_slide.worked_item.formula_h", min(0.42, formula_h))
        if display_line_count >= 7:
            formula_size = style_value(design, "problem_slide.worked_item.dense_formula_font_size", min(16, formula_size))
            formula_h = style_value(design, "problem_slide.worked_item.dense_formula_h", min(0.34, formula_h))
    group = slide.shapes.add_group_shape()
    item_label = item_group_label(item, idx)
    group.name = f"reveal-item-{idx}-{item_label}" if reveal else f"visible-item-{idx}-{item_label}"
    number_shape = add_text(
        group.shapes,
        number_x,
        y,
        style_value(design, "problem_slide.item.number_w", 0.70),
        style_value(design, "problem_slide.item.number_h", 0.38),
        f"({item.number})" if item.number is not None else f"({idx + 1})",
        size=style_value(design, "problem_slide.item.number_font_size", 21),
        bold=True,
    )
    number_shape.name = f"{item_label}-number"
    segment_rows = display_segment_rows(item)
    use_segment_rows = display_segments_are_structured(item, segment_rows)
    lines = [] if use_segment_rows else split_worked_expression_lines(item)
    line_gap = style_value(design, "problem_slide.worked_item.line_gap", 0.40) if multiline else formula_h
    if multiline and display_line_count >= 7:
        line_gap = style_value(design, "problem_slide.worked_item.dense_line_gap", min(0.31, line_gap))
    if multiline:
        line_gap = worked_item_line_gap(item, design, line_gap)
    formula_y_offset = style_value(design, "problem_slide.item.formula_y_offset", -0.02)
    if item.layout_shapes:
        bracket_group = group.shapes.add_group_shape()
        bracket_group.name = f"{item_label}-bracket"
        rendered_count = add_item_layout_shapes(
            bracket_group.shapes,
            item,
            number_x=number_x,
            formula_x=formula_x,
            y=y,
            line_gap=line_gap,
            formula_y_offset=formula_y_offset,
            design=design,
        )
        if rendered_count == 0:
            bracket_group.name = f"{item_label}-bracket-empty"
    for line_idx, segments in enumerate(segment_rows if use_segment_rows else []):
        line_y = y + formula_y_offset + line_idx * line_gap
        line_group = group.shapes.add_group_shape()
        line_group.name = f"{item_label}-line-{line_idx + 1}"
        add_display_segment_row(
            line_group.shapes,
            formula_x,
            line_y,
            formula_w,
            formula_h,
            segments,
            design=design,
            size=formula_size,
        )
    for line_idx, expr in enumerate(lines):
        line_y = y + formula_y_offset + line_idx * line_gap
        line_group = group.shapes.add_group_shape()
        line_group.name = f"{item_label}-line-{line_idx + 1}"
        add_math_row(
            line_group.shapes,
            formula_x,
            line_y,
            formula_w,
            formula_h,
            expr,
            design=design,
            size=formula_size,
            font=style_value(design, "fonts.math", MATH_FONT),
        )
    return group.name if reveal else None


def vertical_layout_override_key(block: PracticeBlock, items: list[PracticeItem]) -> str | None:
    key = f"page{block.page}_practice{block.practice_no}"
    if block.page == 14 and block.practice_no == "1":
        suffix = "three_items" if len(items) >= 3 else "two_items"
        return f"{key}_{suffix}"
    if block.page in (10, 11) and block.practice_no in ("1", "2"):
        return key
    return None


def add_vertical_items(
    slide,
    block: PracticeBlock,
    items: list[PracticeItem],
    design: dict[str, Any] | None,
    *,
    start_y: float | None = None,
) -> None:
    if start_y is None:
        start_y = style_value(design, "problem_slide.item.start_y", 2.749)
    row_gap = style_value(design, "problem_slide.item.row_gap", 1.12)
    override_key = vertical_layout_override_key(block, items)
    if override_key:
        row_gap = style_value(design, f"problem_slide.layout_overrides.{override_key}.row_gap", row_gap)
    max_segment_rows = max((len(display_segment_rows(item)) for item in items), default=1)
    if max_segment_rows >= 2:
        worked_line_gap = style_value(design, "problem_slide.worked_item.line_gap", 0.38)
        row_gap = max(row_gap, max_segment_rows * worked_line_gap + 0.22)
    for idx, item in enumerate(items):
        add_item_group(slide, item, idx, reveal=idx > 0, design=design, y=start_y + idx * row_gap)


def add_two_column_items(
    slide,
    items: list[PracticeItem],
    design: dict[str, Any] | None,
    *,
    start_y: float | None = None,
) -> None:
    ordered_items = display_order(items, "two_column_grid")
    has_multiline = any(is_multiline_item(item) for item in ordered_items)
    has_explicit_grid_positions = any(item.source_box for item in ordered_items) or any(
        item.column_index == 1 for item in ordered_items
    )
    position_rows = [item.row_index for item in ordered_items] if has_explicit_grid_positions else []
    row_offset = min(position_rows) if position_rows else 0
    left_number_x = style_value(design, "problem_slide.two_column_grid.left_number_x", 0.35)
    left_formula_x = style_value(design, "problem_slide.two_column_grid.left_formula_x", 1.05)
    right_number_x = style_value(design, "problem_slide.two_column_grid.right_number_x", 5.35)
    right_formula_x = style_value(design, "problem_slide.two_column_grid.right_formula_x", 6.05)
    if start_y is None:
        start_y = style_value(design, "problem_slide.two_column_grid.start_y", 3.371)
    row_gap = style_value(design, "problem_slide.two_column_grid.row_gap", 1.75)
    if has_multiline:
        start_y = style_value(design, "problem_slide.worked_grid.start_y", start_y)
        row_gap = style_value(design, "problem_slide.worked_grid.row_gap", 2.35)
    formula_w = style_value(design, "problem_slide.two_column_grid.formula_w", 3.55)
    formula_h = style_value(design, "problem_slide.two_column_grid.formula_h", 0.62)
    for idx, item in enumerate(ordered_items):
        if has_explicit_grid_positions:
            column = item.column_index if item.column_index in (0, 1) else idx % 2
            row = item.row_index - row_offset
        else:
            column = idx % 2
            row = idx // 2
        add_item_group(
            slide,
            item,
            idx,
            reveal=idx > 0,
            design=design,
            number_x=right_number_x if column else left_number_x,
            formula_x=right_formula_x if column else left_formula_x,
            y=start_y + row * row_gap,
            formula_w=formula_w,
            formula_h=formula_h,
        )


def add_problem_slide(prs: Presentation, block: PracticeBlock, items: list[PracticeItem], continuation: int, design: dict[str, Any] | None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_practice_header(slide, block.practice_no, design)
    prompt = " ".join(block.prompt.split())
    prompt_y = style_value(design, "problem_slide.prompt.y", 0.993)
    prompt_h = style_value(design, "problem_slide.prompt.h", 1.254)
    is_long_prompt = len(prompt) >= style_value(design, "problem_slide.prompt.long_text_threshold", 64)
    if is_long_prompt:
        prompt_h = style_value(design, "problem_slide.prompt.long_h", prompt_h)
    prompt_gap = style_value(design, "problem_slide.item.prompt_gap", 0.502)
    if is_long_prompt:
        prompt_gap = style_value(design, "problem_slide.item.long_prompt_gap", prompt_gap)
    item_start_y = prompt_y + prompt_h + prompt_gap
    prompt_w = style_value(design, "problem_slide.prompt.w", 9.428)
    override_key = vertical_layout_override_key(block, items)
    if override_key:
        prompt_w = style_value(design, f"problem_slide.layout_overrides.{override_key}.prompt_w", prompt_w)
    prompt_shape = add_text(
        slide.shapes,
        style_value(design, "problem_slide.prompt.x", 0.224),
        prompt_y,
        prompt_w,
        prompt_h,
        prompt,
        size=style_value(design, "problem_slide.prompt.font_size", 24),
        bold=True,
        line_spacing=style_value(design, "problem_slide.prompt.line_spacing", 1.5),
    )
    if block.layout_type == "two_column_grid":
        add_two_column_items(slide, items, design, start_y=item_start_y)
    else:
        add_vertical_items(slide, block, items, design, start_y=item_start_y)


def build_presentation(blocks: list[PracticeBlock], output: Path, design: dict[str, Any] | None = None) -> dict:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    last_concept_key: tuple[str, str] | None = None
    slide_trace = []
    title_slide_numbers: list[int] = []

    for block in blocks:
        concept_key = (block.concept_no, block.concept_title)
        if concept_key != last_concept_key:
            add_title_slide(prs, block.concept_no, block.concept_title, design)
            title_slide_numbers.append(len(prs.slides))
            last_concept_key = concept_key
        chunks = split_block_items(block)
        for continuation, chunk in enumerate(chunks):
            add_problem_slide(prs, block, chunk, continuation, design)
            slide_trace.append(
                {
                    "slide_number": len(prs.slides),
                    "slide_type": "problem",
                    "page": block.page,
                    "practice": block.practice_no,
                    "concept_no": block.concept_no,
                    "concept_title": block.concept_title,
                    "prompt": block.prompt,
                    "layout_type": block.layout_type,
                    "chunk_index": continuation + 1,
                    "chunk_count": len(chunks),
                    "items": [item.raw_text for item in chunk],
                    "item_numbers": [item.number for item in chunk],
                    "visible_by_default_item_numbers": [chunk[0].number] if chunk else [],
                    "click_reveal_item_numbers": [item.number for item in chunk[1:]],
                    "item_inventory": [item.to_dict() for item in chunk],
                }
            )

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    style_summary = apply_ooxml_style_parts(output, design, title_slide_numbers=title_slide_numbers)
    inject_group_reveals(output)
    return {"slide_trace": slide_trace, "slide_count": len(prs.slides), "ooxml_style": style_summary}


def p_tag(name: str) -> str:
    return f"{{{P_NS}}}{name}"


def timing_branch(target_id: str, start_id: int) -> ET.Element:
    outer_id, inner_id, effect_id, set_id = start_id, start_id + 1, start_id + 2, start_id + 3
    par1 = ET.Element(p_tag("par"))
    ctn1 = ET.SubElement(par1, p_tag("cTn"), id=str(outer_id), fill="hold")
    st1 = ET.SubElement(ctn1, p_tag("stCondLst"))
    ET.SubElement(st1, p_tag("cond"), delay="indefinite")
    child1 = ET.SubElement(ctn1, p_tag("childTnLst"))
    par2 = ET.SubElement(child1, p_tag("par"))
    ctn2 = ET.SubElement(par2, p_tag("cTn"), id=str(inner_id), fill="hold")
    st2 = ET.SubElement(ctn2, p_tag("stCondLst"))
    ET.SubElement(st2, p_tag("cond"), delay="0")
    child2 = ET.SubElement(ctn2, p_tag("childTnLst"))
    par3 = ET.SubElement(child2, p_tag("par"))
    ctn3 = ET.SubElement(
        par3,
        p_tag("cTn"),
        id=str(effect_id),
        presetID="1",
        presetClass="entr",
        presetSubtype="0",
        fill="hold",
        nodeType="clickEffect",
    )
    st3 = ET.SubElement(ctn3, p_tag("stCondLst"))
    ET.SubElement(st3, p_tag("cond"), delay="0")
    child3 = ET.SubElement(ctn3, p_tag("childTnLst"))
    set_node = ET.SubElement(child3, p_tag("set"))
    c_bhvr = ET.SubElement(set_node, p_tag("cBhvr"))
    ctn4 = ET.SubElement(c_bhvr, p_tag("cTn"), id=str(set_id), dur="1", fill="hold")
    st4 = ET.SubElement(ctn4, p_tag("stCondLst"))
    ET.SubElement(st4, p_tag("cond"), delay="0")
    tgt_el = ET.SubElement(c_bhvr, p_tag("tgtEl"))
    ET.SubElement(tgt_el, p_tag("spTgt"), spid=target_id)
    attrs = ET.SubElement(c_bhvr, p_tag("attrNameLst"))
    ET.SubElement(attrs, p_tag("attrName")).text = "style.visibility"
    to = ET.SubElement(set_node, p_tag("to"))
    ET.SubElement(to, p_tag("strVal"), val="visible")
    return par1


def build_timing(target_ids: list[str]) -> ET.Element:
    timing = ET.Element(p_tag("timing"))
    tn_lst = ET.SubElement(timing, p_tag("tnLst"))
    root_par = ET.SubElement(tn_lst, p_tag("par"))
    root_ctn = ET.SubElement(root_par, p_tag("cTn"), id="1", dur="indefinite", restart="never", nodeType="tmRoot")
    root_child = ET.SubElement(root_ctn, p_tag("childTnLst"))
    seq = ET.SubElement(root_child, p_tag("seq"), concurrent="1", nextAc="seek")
    seq_ctn = ET.SubElement(seq, p_tag("cTn"), id="2", dur="indefinite", nodeType="mainSeq")
    child = ET.SubElement(seq_ctn, p_tag("childTnLst"))
    next_id = 3
    for target_id in target_ids:
        child.append(timing_branch(target_id, next_id))
        next_id += 4
    for tag, evt in (("prevCondLst", "onPrev"), ("nextCondLst", "onNext")):
        cond_lst = ET.SubElement(seq, p_tag(tag))
        cond = ET.SubElement(cond_lst, p_tag("cond"), evt=evt, delay="0")
        tgt = ET.SubElement(cond, p_tag("tgtEl"))
        ET.SubElement(tgt, p_tag("sldTgt"))
    return timing


def inject_group_reveals(pptx_path: Path) -> None:
    with tempfile.TemporaryDirectory() as tmp:
        tmp_path = Path(tmp)
        with zipfile.ZipFile(pptx_path) as src:
            src.extractall(tmp_path)

        slide_dir = tmp_path / "ppt" / "slides"
        for slide_xml in sorted(slide_dir.glob("slide*.xml"), key=lambda path: int(re.search(r"\d+", path.stem).group(0))):
            tree = ET.parse(slide_xml)
            root = tree.getroot()
            target_ids: list[str] = []
            for grp in root.findall(f".//{p_tag('grpSp')}"):
                c_nv_pr = grp.find(f"./{p_tag('nvGrpSpPr')}/{p_tag('cNvPr')}")
                if c_nv_pr is None:
                    continue
                name = c_nv_pr.get("name") or ""
                if name.startswith("reveal-item-"):
                    c_nv_pr.attrib.pop("hidden", None)
                    target_id = c_nv_pr.get("id")
                    if target_id:
                        target_ids.append(target_id)
            if not target_ids:
                continue
            existing = root.find(p_tag("timing"))
            if existing is not None:
                root.remove(existing)
            root.append(build_timing(target_ids))
            tree.write(slide_xml, encoding="UTF-8", xml_declaration=True)

        rebuilt = pptx_path.with_suffix(".tmp.pptx")
        with zipfile.ZipFile(rebuilt, "w", zipfile.ZIP_DEFLATED) as dst:
            for path in tmp_path.rglob("*"):
                if path.is_file():
                    dst.write(path, path.relative_to(tmp_path).as_posix())
        shutil.move(str(rebuilt), pptx_path)
